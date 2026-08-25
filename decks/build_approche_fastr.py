#!/usr/bin/env python3
"""Build a natively editable PPTX for the "L'approche FASTR" advocacy deck.

Rebuilds the patched-up source deck (`Approche FASTR.pptx`) on the FASTR design
system used by the plaidoyer deck: real editable text boxes and shapes, the
brand palette, flat accent-bar layouts (no shaded/bordered boxes), section
dividers, page-numbered footers, and a fade transition on every slide.
Charts/screenshots are reused from resources/screenshots/approche_fastr/.

Usage:  python3 decks/build_approche_fastr.py
Deps:   pip install python-pptx pillow
"""
from __future__ import annotations

import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parent.parent
RES = REPO / "resources"
IMG = RES / "screenshots" / "approche_fastr"
FONT = "Poppins"

DEEP_GREEN = RGBColor(0x09, 0x54, 0x4F)
GREEN = RGBColor(0x1F, 0x9A, 0x9C)
LIME = RGBColor(0xD0, 0xCB, 0x17)
NAVY = RGBColor(0x21, 0x56, 0x8C)
GOLD = RGBColor(0xD8, 0xA8, 0x22)
INK = RGBColor(0x1A, 0x1F, 0x1E)
INK2 = RGBColor(0x5A, 0x65, 0x62)
INK3 = RGBColor(0x97, 0xA0, 0x9D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HAIRLINE = RGBColor(0xD5, 0xDA, 0xD8)
RED = RGBColor(0xC5, 0x30, 0x30)

EMU_IN = 914400
SW, SH = 13.333, 7.5
LEFT = 0.62
CW = SW - 2 * LEFT
STEP_COLORS = [NAVY, DEEP_GREEN, GREEN, GOLD]
TINT = {NAVY: RGBColor(0xED, 0xF1, 0xF7), DEEP_GREEN: RGBColor(0xE8, 0xF1, 0xF0),
        GREEN: RGBColor(0xE9, 0xF4, 0xF4), GOLD: RGBColor(0xF9, 0xF3, 0xE4)}
CENTER = PP_ALIGN.CENTER
MID = MSO_ANCHOR.MIDDLE


# ---- helpers -------------------------------------------------------------
def _runs(text):
    out, i = [], 0
    for m in re.finditer(r"\*\*(.+?)\*\*", text):
        if m.start() > i:
            out.append((text[i:m.start()], False))
        out.append((m.group(1), True))
        i = m.end()
    if i < len(text):
        out.append((text[i:], False))
    return out or [(text, False)]


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tb, tf


def _para(tf, text, *, size, color, first=False, bold=False, italic=False,
          align=PP_ALIGN.LEFT, bullet=False, bullet_color=LIME, space=8, line=1.1):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    p.space_before = Pt(0)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if bullet:
        r = p.add_run()
        r.text = "•  "
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.bold = True
        r.font.color.rgb = bullet_color
    for seg, b in _runs(text):
        r = p.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.bold = bool(b or bold)
        r.font.italic = italic
        r.font.color.rgb = color
    return p


def _rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    # Drop the inherited theme style — it carries a drop shadow (the "shading"
    # that otherwise appears under every rule/card). We set fill/line explicitly.
    el = sp._element
    st = el.find(qn("p:style"))
    if st is not None:
        el.remove(st)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    return sp


def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _fullbleed(slide, img):
    slide.shapes.add_picture(str(img), 0, 0, Inches(SW), Inches(SH))


def _fit(img_path, max_w, max_h):
    with Image.open(img_path) as im:
        w, h = im.size
    r, R = w / h, max_w / max_h
    return (max_w, max_w / r) if r > R else (max_h * r, max_h)


def _blocks(tf, blocks, base=14):
    first = True
    for b in blocks:
        bl = b.get("bullet", False)
        _para(tf, b["t"], size=b.get("size", base), color=b.get("color", INK), first=first,
              bullet=bl, bold=b.get("bold", False), italic=b.get("italic", False),
              space=b.get("space", 10 if not bl else 6), line=b.get("line", 1.2))
        first = False


def _title(slide, text, *, y=0.5):
    lines = max(1, -(-len(text) // 52))
    th = 0.6 * lines
    _, tf = _tb(slide, LEFT, y, CW, th + 0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _para(tf, text, size=30, color=DEEP_GREEN, first=True, bold=True, space=0, line=1.02)
    ry = y + th + 0.06
    rw = min(CW, 0.6 + 0.14 * len(text)) if lines == 1 else min(CW, 8.5)
    _rect(slide, LEFT, ry, rw, 0.045, fill=GREEN)
    return ry + 0.3


def _callout(slide, text):
    """Key-takeaway line — left accent bar, no background shading."""
    y = SH - 1.02
    _rect(slide, LEFT, y, 0.06, 0.5, fill=DEEP_GREEN)
    _, tf = _tb(slide, LEFT + 0.24, y, CW - 0.35, 0.5)
    tf.vertical_anchor = MID
    _para(tf, text, size=12.5, color=DEEP_GREEN, first=True, italic=True, bold=True, space=0, line=1.15)


def _footer(slide, n, total):
    _rect(slide, LEFT, SH - 0.44, CW, 0.012, fill=RGBColor(0xE4, 0xE7, 0xE5))
    _, lf = _tb(slide, LEFT, SH - 0.38, CW - 1.2, 0.3)
    _para(lf, "FASTR  ·  L'APPROCHE FASTR", size=8, color=INK3, first=True, bold=True, space=0)
    _, pf = _tb(slide, SW - LEFT - 1.0, SH - 0.38, 1.0, 0.3)
    _para(pf, f"{n} / {total}", size=8, color=INK3, first=True, bold=True, align=PP_ALIGN.RIGHT, space=0)


def _fade(slide, spd="med"):
    sld = slide._element
    for t in sld.findall(qn("p:transition")):
        sld.remove(t)
    sld.append(parse_xml(
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'spd="{spd}"><p:fade/></p:transition>'))


def _new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---- archetypes ----------------------------------------------------------
def s_cover(prs, c):
    s = _new(prs)
    _fullbleed(s, RES / "backgrounds" / "cover_slide_clean.png")
    s.shapes.add_picture(str(RES / "logos" / "GFF_Logo_trimmed.png"), Inches(0.8), Inches(0.55), height=Inches(0.42))
    _, tf = _tb(s, 0.8, 2.7, 11.2, 2.4)
    tf.vertical_anchor = MID
    _para(tf, c["title"], size=44, color=WHITE, first=True, bold=True, space=10, line=1.03)
    _para(tf, c["subtitle"], size=18, color=RGBColor(0xEC, 0xF0, 0xEF), space=0, line=1.2)
    y = SH - 1.05
    s.shapes.add_picture(str(RES / "logos" / "FASTR_White_Horiz.png"), Inches(0.8), Inches(y), height=Inches(0.5))
    s.shapes.add_picture(str(RES / "logos" / "usefuldata600w.png"), Inches(3.0), Inches(y + 0.07), height=Inches(0.34))


def s_section(prs, c):
    s = _new(prs)
    _fullbleed(s, RES / "backgrounds" / "section_slide.png")
    if c.get("kicker"):
        _, kf = _tb(s, 1.4, 2.45, SW - 2.8, 0.5)
        _para(kf, c["kicker"].upper(), size=13, color=LIME, first=True, bold=True, align=CENTER, space=0)
    _, tf = _tb(s, 1.2, 0, SW - 2.4, SH)
    tf.vertical_anchor = MID
    _para(tf, c["title"], size=32, color=WHITE, first=True, bold=True, align=CENTER, space=0, line=1.1)


def s_statement(prs, c):
    """A single strong statement on white — for the objective."""
    s = _new(prs)
    _, tf = _tb(s, 1.4, 0, SW - 2.8, SH)
    tf.vertical_anchor = MID
    first = True
    if c.get("kicker"):
        _para(tf, c["kicker"].upper(), size=13, color=GREEN, first=True, bold=True, align=CENTER, space=16)
        first = False
    _para(tf, c["text"], size=27, color=DEEP_GREEN, first=first, bold=True, align=CENTER, space=0, line=1.24)


def s_content(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    _, tf = _tb(s, LEFT, by, c.get("tw", CW), SH - by - 0.6)
    _blocks(tf, c["blocks"], base=16)


def s_cycle(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    n = len(c["nodes"])
    arrow = 0.5
    bw = (CW - (n - 1) * arrow) / n
    bh, y0 = 3.0, by + 0.5
    for i, (label, desc) in enumerate(c["nodes"]):
        x = LEFT + i * (bw + arrow)
        col = STEP_COLORS[i % 4]
        _rect(s, x, y0, bw, bh, fill=TINT[col])
        _rect(s, x, y0, bw, 0.09, fill=col)
        _, tf = _tb(s, x + 0.24, y0 + 0.34, bw - 0.48, bh - 0.5)
        _para(tf, label, size=18, color=col, first=True, bold=True, align=CENTER, space=10, line=1.05)
        _para(tf, desc, size=12.5, color=INK2, align=CENTER, space=0, line=1.22)
        if i < n - 1:
            _, af = _tb(s, x + bw, y0, arrow, bh)
            af.vertical_anchor = MID
            _para(af, "→", size=22, color=INK3, first=True, align=CENTER)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_pillars(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    n = 4
    gap = 0.45
    cw = (CW - (n - 1) * gap) / n
    ch, y0 = 3.5, by + 0.15
    for i, p in enumerate(c["pillars"]):
        x = LEFT + i * (cw + gap)
        col = STEP_COLORS[i % 4]
        _rect(s, x, y0, cw, ch, fill=TINT[col])
        _rect(s, x, y0, cw, 0.09, fill=col)
        _, tf = _tb(s, x + 0.24, y0 + 0.26, cw - 0.44, ch - 0.4)
        _para(tf, f"0{i+1}", size=16, color=col, first=True, bold=True, space=7)
        _para(tf, p["head"], size=14.5, color=INK, bold=True, space=9, line=1.1)
        _para(tf, p["desc"], size=11.5, color=INK2, space=0, line=1.24)
    yb = y0 + ch + 0.28
    _rect(s, LEFT, yb, CW, 0.68, fill=DEEP_GREEN)
    _, bf = _tb(s, LEFT + 0.35, yb, CW - 0.7, 0.68)
    bf.vertical_anchor = MID
    p = _para(bf, "", size=14, color=WHITE, first=True, space=0, align=CENTER)
    r = p.add_run(); r.text = "Transversal    "; r.font.size = Pt(13); r.font.name = FONT; r.font.bold = True; r.font.color.rgb = LIME
    r = p.add_run(); r.text = c["crosscutting"]; r.font.size = Pt(14); r.font.name = FONT; r.font.bold = True; r.font.color.rgb = WHITE


def s_numbered(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    n = len(c["steps"])
    gap = 0.45
    bw = (CW - (n - 1) * gap) / n
    bh, y0 = 3.3, by + 0.3
    for i, (head, desc) in enumerate(c["steps"]):
        x = LEFT + i * (bw + gap)
        col = STEP_COLORS[i % 4]
        _rect(s, x, y0, bw, bh, fill=TINT[col])
        _rect(s, x, y0, bw, 0.09, fill=col)
        _, nf = _tb(s, x + 0.24, y0 + 0.3, bw - 0.48, 0.9)
        _para(nf, f"0{i+1}", size=38, color=col, first=True, bold=True, space=0)
        _, tf = _tb(s, x + 0.24, y0 + 1.28, bw - 0.48, bh - 1.45)
        _para(tf, head, size=15, color=INK, first=True, bold=True, space=8, line=1.05)
        _para(tf, desc, size=12, color=INK2, space=0, line=1.24)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_image_text(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    img = IMG / c["image"]
    tw = c.get("tw", 4.5)
    iw_box = CW - tw - 0.45
    ih_box = SH - by - 0.6
    w, h = _fit(img, iw_box, ih_box)
    right = c.get("side", "left") == "right"
    ix = (SW - LEFT - w) if right else (LEFT + (iw_box - w) / 2)
    tx = LEFT if right else LEFT + iw_box + 0.45
    s.shapes.add_picture(str(img), Inches(ix), Inches(by + (ih_box - h) / 2), Inches(w), Inches(h))
    _, tf = _tb(s, tx, by, tw, ih_box)
    tf.vertical_anchor = MID
    _blocks(tf, c["blocks"], base=14)


def s_image_full(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    img = IMG / c["image"]
    area = SH - by - 0.6
    w, h = _fit(img, CW, area)
    s.shapes.add_picture(str(img), Inches((SW - w) / 2), Inches(by + (area - h) / 2), Inches(w), Inches(h))


def s_table(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    if c.get("intro"):
        _, itf = _tb(s, LEFT, by, CW, 0.5)
        _para(itf, c["intro"], size=13, color=INK2, first=True, italic=True, space=0, line=1.15)
        by += 0.5
    rows, cols = len(c["rows"]) + 1, len(c["headers"])
    th = min(SH - by - 1.05, 0.42 * rows)
    gt = s.shapes.add_table(rows, cols, Inches(LEFT), Inches(by), Inches(CW), Inches(th)).table
    if c.get("colw"):
        for j, w in enumerate(c["colw"]):
            gt.columns[j].width = Inches(w)
    for j, htxt in enumerate(c["headers"]):
        cell = gt.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = DEEP_GREEN
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.vertical_anchor = MID
        _para(cell.text_frame, htxt, size=11.5, color=WHITE, first=True, bold=True, space=0, line=1.05)
    for i, row in enumerate(c["rows"], 1):
        for j, val in enumerate(row):
            cell = gt.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF4, 0xF7, 0xF6)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MID
            bold = (j == 0 and c.get("bold_first"))
            _para(cell.text_frame, val, size=c.get("fs", 11), color=DEEP_GREEN if bold else INK,
                  first=True, bold=bold, space=0, line=1.08)
    if c.get("note"):
        _, nf = _tb(s, LEFT, SH - 1.02, CW, 0.55)
        _para(nf, c["note"], size=9.5, color=INK3, first=True, italic=True, space=0, line=1.12)


def s_two_col(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    widths = c.get("widths", [1, 1])
    tot = sum(widths)
    x = LEFT
    for i, col in enumerate(c["cols"]):
        w = (CW - 0.7) * widths[i] / tot
        _, tf = _tb(s, x, by, w, SH - by - 0.7)
        _para(tf, col["head"], size=16, color=DEEP_GREEN, first=True, bold=True, space=9, line=1.05)
        for it in col["items"]:
            _para(tf, it, size=col.get("size", 14), color=INK, bullet=col.get("bullet", True), space=5, line=1.16)
        if col.get("note"):
            _para(tf, col["note"], size=12, color=DEEP_GREEN, italic=True, bold=True, space=0, line=1.12)
        x += w + 0.7


def s_demo(prs, c):
    s = _new(prs)
    _, ef = _tb(s, LEFT, 0.55, CW, 0.4)
    _para(ef, c.get("eyebrow", "DÉMONSTRATION EN DIRECT").upper(), size=13, color=GREEN, first=True, bold=True, space=0)
    _, tf0 = _tb(s, LEFT, 1.0, CW, 1.0)
    _para(tf0, c["title"], size=32, color=DEEP_GREEN, first=True, bold=True, space=0)
    _rect(s, LEFT, 1.78, min(CW, 0.6 + 0.14 * len(c["title"])), 0.045, fill=GREEN)
    _, tf = _tb(s, LEFT, 2.35, 8.6, SH - 3.1)
    _para(tf, c["lead"], size=16, color=INK, first=True, space=18, line=1.2)
    for it in c["cues"]:
        _para(tf, it, size=15.5, color=INK, bullet=True, space=11, line=1.16)


def s_ai_flow(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    _, pf = _tb(s, LEFT, by, CW, 0.6)
    _para(pf, c["principle"], size=16, color=INK, first=True, space=0, line=1.2)
    ry = by + 0.95
    rh, rgap = 1.05, 0.55
    for r in c["rows"]:
        col = r["color"]
        steps = r["steps"]
        label_w, arrow = 3.0, 0.34
        sw = (CW - label_w - arrow * len(steps)) / len(steps)
        x = LEFT
        _rect(s, x, ry, label_w, rh, fill=col)
        _, lf = _tb(s, x + 0.18, ry, label_w - 0.36, rh)
        lf.vertical_anchor = MID
        _para(lf, r["label"], size=14, color=WHITE, first=True, bold=True, align=CENTER, line=1.08)
        x += label_w
        for st in steps:
            _, af = _tb(s, x, ry, arrow, rh)
            af.vertical_anchor = MID
            _para(af, "→", size=18, color=INK3, first=True, align=CENTER)
            x += arrow
            _rect(s, x, ry, sw, rh, fill=TINT[col])
            _rect(s, x, ry, sw, 0.07, fill=col)
            _, cf = _tb(s, x + 0.15, ry, sw - 0.3, rh)
            cf.vertical_anchor = MID
            _para(cf, st, size=13, color=INK, first=True, align=CENTER, space=0, line=1.12)
            x += sw
        ry += rh + rgap


def s_contrast(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    half = (CW - 0.6) / 2
    ph, y0 = 2.85, by + 0.1
    for i, col in enumerate(c["cols"]):
        x = LEFT + i * (half + 0.6)
        acc = col["accent"]
        _rect(s, x, y0 + 0.02, 0.05, ph - 0.04, fill=acc)
        _, tf = _tb(s, x + 0.26, y0, half - 0.26, ph)
        _para(tf, col["head"], size=16.5, color=acc, first=True, bold=True, space=10, line=1.08)
        for it in col["items"]:
            _para(tf, it, size=13.5, color=INK, bullet=True, space=7, line=1.2)
    for i, (tag, desc, acc) in enumerate(c.get("bottom", [])):
        cw2 = (CW - 0.5) / 2
        cy = y0 + ph + 0.35
        x = LEFT + i * (cw2 + 0.5)
        _rect(s, x, cy, cw2, 0.92, fill=None, line=acc, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(s, x, cy, 0.06, 0.92, fill=acc)
        _, tf = _tb(s, x + 0.28, cy, cw2 - 0.45, 0.92)
        tf.vertical_anchor = MID
        _para(tf, tag, size=14, color=acc, first=True, bold=True, space=3)
        _para(tf, desc, size=12.5, color=INK, space=0, line=1.12)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_closing(prs, c):
    s = _new(prs)
    _bg(s, DEEP_GREEN)
    _, tf = _tb(s, 0.9, 0.9, 11.4, 1.2)
    _para(tf, c["title"], size=40, color=WHITE, first=True, bold=True, space=0)
    _rect(s, 0.95, 2.05, 1.6, 0.05, fill=LIME)
    _, tf2 = _tb(s, 0.9, 2.35, 11.4, 1.2)
    _para(tf2, c["tagline"], size=22, color=WHITE, first=True, bold=True, space=0, line=1.12)
    cy = 4.65  # shared vertical center for the logo row
    s.shapes.add_picture(str(RES / "logos" / "FASTR_White_Horiz.png"), Inches(0.95), Inches(cy - 0.25), height=Inches(0.5))
    s.shapes.add_picture(str(RES / "logos" / "usefuldata600w.png"), Inches(3.25), Inches(cy - 0.17), height=Inches(0.34))
    s.shapes.add_picture(str(RES / "logos" / "GFF_Logo_trimmed.png"), Inches(6.15), Inches(cy - 0.15), height=Inches(0.3))


BUILDERS = {"cover": s_cover, "section": s_section, "statement": s_statement, "content": s_content,
            "cycle": s_cycle, "pillars": s_pillars, "numbered": s_numbered, "image_text": s_image_text,
            "image_full": s_image_full, "table": s_table, "two_col": s_two_col, "demo": s_demo,
            "ai_flow": s_ai_flow, "contrast": s_contrast, "closing": s_closing}


CONTENT = [
    {"type": "cover", "title": "L'approche FASTR",
     "subtitle": "Des analyses à cycle rapide pour accélérer les progrès de la SRMNEA-N"},

    {"type": "statement", "kicker": "Objectif",
     "text": "Les analyses à cycle rapide accélèrent les progrès de la SRMNEA-N en renforçant la disponibilité et l'utilisation systématique et opportune des données pour la prise de décisions stratégiques."},

    {"type": "content", "title": "Qu'est-ce que FASTR ?", "tw": 10.8, "blocks": [
        {"t": "Le GFF soutient les efforts menés par les pays pour améliorer la génération et l'utilisation, en temps utile, des données pour la prise de décision — afin de renforcer les systèmes de soins de santé primaires (SSP) et d'améliorer les résultats en matière de SRMNEA-N.", "size": 16, "space": 13},
        {"t": "Cet ensemble d'initiatives et d'appui technique est désigné sous le nom de **Frequent Assessments and System Tools for Resilience (FASTR)**.", "size": 16, "space": 13},
        {"t": "FASTR catalyse des cycles continus d'**analyse, apprentissage, renforcement et action** pour ancrer l'utilisation systématique des données dans la décision.", "size": 16, "space": 0},
    ]},

    {"type": "cycle", "title": "Un cycle continu", "nodes": [
        ("Analyser", "Analyser en temps utile les données de SSP et de SRMNEA-N"),
        ("Apprendre", "Apprendre des données pour prioriser les lacunes des systèmes de santé"),
        ("Renforcer", "Renforcer les capacités de S&E et la qualité des données dans le pays"),
        ("Agir", "Traduire les données en actions pour de meilleurs résultats SRMNEA-N"),
    ], "footer": "…et le cycle recommence, trimestre après trimestre."},

    {"type": "content", "title": "Comment y parvenons-nous ?", "tw": 10.8, "blocks": [
        {"t": "FASTR propose un ensemble d'approches **rigoureuses, rapides et peu coûteuses** pour le suivi et l'évaluation des systèmes de soins de santé primaires.", "size": 17},
        {"t": "Elles s'accompagnent d'un **renforcement des capacités** et d'un **soutien à l'utilisation des données**, adaptés aux besoins du pays.", "size": 17},
    ]},

    {"type": "section", "kicker": "Partie 1", "title": "Une approche intégrée : quatre piliers complémentaires"},

    {"type": "pillars", "title": "Quatre piliers complémentaires", "pillars": [
        {"head": "Utilisation des services SRMNEA-N", "desc": "Exploiter les données du SNIS pour suivre l'évolution de la couverture des services — en période de réforme, de choc ou au sein des populations vulnérables — tout en corrigeant les problèmes de qualité des données."},
        {"head": "Enquêtes auprès des formations sanitaires", "desc": "Évaluer les lacunes dans la prestation des soins de santé primaires et comprendre l'impact des réformes ou des chocs sur la performance des systèmes de santé."},
        {"head": "Enquête auprès des ménages et des usagers", "desc": "Fournir un aperçu de l'utilisation des services, du renoncement aux soins, et des perceptions communautaires de la qualité des services."},
        {"head": "Études qualitatives rapides", "desc": "Mieux comprendre les enjeux émergents des systèmes de santé à travers des approches mixtes, rapides et flexibles."},
    ], "crosscutting": "Renforcement des capacités et soutien à l'utilisation des données"},

    {"type": "image_full", "title": "FASTR est opérationnel dans 29 pays partenaires du GFF — et continue de se développer",
     "image": "s07_Picture_6.png"},

    {"type": "section", "kicker": "Partie 2", "title": "Enquêtes rapides auprès des formations sanitaires"},

    {"type": "content", "title": "Quelle est l'approche FASTR ?", "tw": 10.8, "blocks": [
        {"t": "Nous concevons et mettons en œuvre des **enquêtes téléphoniques trimestrielles** auprès des formations sanitaires primaires.", "size": 16},
        {"t": "Elles produisent des analyses régulières de la performance des soins de santé primaires, de la disponibilité et de la capacité opérationnelle des services, sur des indicateurs prioritaires liés au plan stratégique SRMNIA-N.", "size": 16},
        {"t": "Ces résultats informent la prise de décision et des actions concrètes pour renforcer le système de santé.", "size": 16},
    ]},

    {"type": "numbered", "title": "Ce que l'enquête permet", "steps": [
        ("Contraintes de l'offre", "Comprendre les contraintes liées à l'offre de services dans les établissements de SSP."),
        ("Mise en œuvre des réformes", "Mesurer la mise en œuvre des réformes sur le terrain."),
        ("Effet des chocs", "Évaluer l'effet des chocs externes sur les systèmes de santé."),
        ("Gestion adaptative", "Améliorer la rapidité et la pertinence des enquêtes comme outil de gestion adaptative."),
    ]},

    {"type": "two_col", "title": "Un outil structuré et adaptable", "widths": [1, 1.15], "cols": [
        {"head": "Structure", "bullet": True, "size": 15, "items": [
            "**10 modules** d'enquête basés sur le cadre **PHCMFI de l'OMS**",
            "Questions inspirées des outils **HHFA/SARA, SPA et SDI**",
        ], "note": "Outil d'enquête 100 % adaptable par les pays."},
        {"head": "Modules de l'outil d'enquête", "size": 14, "items": [
            "Chocs externes", "Résilience aux chocs", "Offre de services", "Infrastructure",
            "Financement", "Ressources humaines", "Fournitures médicales",
            "Direction et coordination", "Participation communautaire",
            "Amélioration de la qualité des soins"]},
    ]},

    {"type": "table", "title": "Conception de l'enquête — exemple du Burkina Faso",
     "intro": "Modalités de mise en œuvre et méthodologie définies par les pays.",
     "headers": ["Paramètre", "Détail"], "colw": [3.0, 9.09], "bold_first": True, "fs": 11.5, "rows": [
        ["Modalité", "Enquête téléphonique auprès des représentants des formations sanitaires."],
        ["Durée", "30–45 minutes par appel."],
        ["Fréquence", "Trimestrielle. L'outil est complété après 4 cycles d'enquête."],
        ["Échantillon", "Panel représentatif de formations sanitaires primaires. Les mêmes formations sont suivies au cours de l'année."],
        ["Échantillonnage", "Représentativité nationale ou infranationale selon les besoins du pays (≈ 200–800 formations par pays)."],
        ["Mise en œuvre", "Gouvernement, partenaire académique/technique ou cabinet privé, selon le contexte et la source de financement."],
        ["Adaptation rapide", "Révision éventuelle de l'outil à chaque trimestre selon les résultats et les questions émergentes."],
     ], "note": "Une assistance technique est disponible pour renforcer les compétences des pays dans la réalisation d'enquêtes téléphoniques rapides."},

    {"type": "table", "title": "Un contenu qui s'adapte de cycle en cycle",
     "intro": "Le contenu évolue à chaque cycle, fournissant de nouvelles perspectives et un apprentissage continu.",
     "headers": ["Cycle 1", "Cycle 2", "Cycle 3", "Cycle 4"], "fs": 10.5, "rows": [
        ["Caractéristiques FOSA", "Caractéristiques FOSA*", "Caractéristiques FOSA*", "Caractéristiques FOSA*"],
        ["—", "Contenu adaptatif**", "Contenu adaptatif**", "Contenu adaptatif**"],
        ["Chocs et résilience", "Chocs et résilience", "Chocs et résilience", "Chocs et résilience"],
        ["Fournitures médicales", "Fournitures médicales", "Fournitures médicales", "Fournitures médicales"],
        ["Infrastructure", "Infrastructure", "Infrastructure", "Infrastructure"],
        ["Offre de services", "Financement", "Amélioration qualité des soins", "Ressources humaines"],
        ["—", "Engagement communautaire", "—", "Leadership et coordination"],
     ], "note": "Modules trimestriels et annuels. *Formations de remplacement uniquement. **Questions supplémentaires identifiées lors de l'adaptation au contexte national ou entre les cycles."},

    {"type": "image_full", "title": "Exemples de résultats — enquête rapide auprès des formations sanitaires",
     "image": "s14_Picture_7.png"},

    {"type": "section", "kicker": "Partie 3", "title": "Analyse de l'utilisation des services SRMNEA-N"},

    {"type": "content", "title": "Quelle est l'approche FASTR ?", "tw": 10.8, "blocks": [
        {"t": "**Analyses trimestrielles des données DHIS2**, axées sur les indicateurs nationaux prioritaires.", "size": 16},
        {"t": "Créer des **outils durables** pour que les parties prenantes génèrent les bonnes analyses et visualisations, au bon moment, sur leurs indicateurs d'intérêt.", "size": 16},
        {"t": "Combiner analyse et visualisation avec le **renforcement des capacités** et le **soutien à l'utilisation des données**, pour la durabilité et l'institutionnalisation.", "size": 16},
    ]},

    {"type": "numbered", "title": "Quatre étapes d'analyse", "steps": [
        ("Qualité des données", "Identifier et corriger les problèmes de qualité pour améliorer la précision de l'analyse."),
        ("Variations de volume", "Quantifier les variations des volumes de services prioritaires au fil du temps."),
        ("Couverture des services", "Comparer les tendances de couverture aux priorités et objectifs du pays."),
        ("Analyses trimestrielles", "Générer des analyses exploitables pour la gestion adaptative."),
    ]},

    {"type": "image_text", "title": "Évaluation et ajustement de la qualité des données",
     "image": "s17_Picture_7.png", "side": "left", "tw": 4.2, "blocks": [
        {"t": "Les évaluations de qualité identifient les problèmes les plus prioritaires.", "space": 12},
        {"t": "Elles déterminent les ajustements analytiques nécessaires.", "space": 12},
        {"t": "L'objectif : éviter que les problèmes de qualité ne deviennent un obstacle à l'analyse et à l'utilisation des données.", "space": 0},
     ]},

    {"type": "image_text", "title": "Surveillance des services de santé essentiels",
     "image": "s18_Picture_4.png", "side": "left", "tw": 4.3, "blocks": [
        {"t": "Nous quantifions les variations du volume des services prioritaires :", "bold": True, "space": 8},
        {"t": "le volume annuel absolu des services sélectionnés", "bullet": True},
        {"t": "la variation en pourcentage du volume mensuel moyen", "bullet": True},
        {"t": "toute année variant de plus de 10 % est signalée", "bullet": True},
        {"t": "des variations importantes peuvent indiquer des perturbations, des réformes ou des anomalies de qualité.", "bullet": True},
     ]},

    {"type": "image_text", "title": "Perturbations et excédents de service",
     "image": "s19_Picture_7.jpg", "side": "left", "tw": 5.2, "blocks": [
        {"t": "Nous évaluons l'impact des perturbations externes et quantifions l'impact des programmes sur l'utilisation des services.", "space": 10},
        {"t": "**Perturbation** — volumes inférieurs aux niveaux prévus : obstacles à l'accès, pénuries de ressources ou défaillances du système.", "space": 10},
        {"t": "**Excédent** — volumes supérieurs aux prévisions : rattrapage, campagnes ou effets de programme.", "space": 0},
     ]},

    {"type": "image_text", "title": "Couverture du service",
     "image": "s20_Picture_4.png", "side": "left", "tw": 4.6, "blocks": [
        {"t": "Des méthodes novatrices calculent et valident les dénominateurs de population, améliorant les estimations de couverture issues du SNIS.", "space": 10},
        {"t": "Là où les données sont précises : identifier les inégalités infranationales et mettre à jour les estimations obsolètes.", "space": 10},
        {"t": "Ailleurs : les tendances observées renseignent tout de même sur la performance globale du système.", "space": 0},
     ]},

    {"type": "section", "kicker": "Partie 4", "title": "La plateforme d'analyse et l'assistant IA"},

    {"type": "image_text", "title": "Extraction des données",
     "image": "s21_Picture_1.png", "side": "right", "tw": 4.6, "blocks": [
        {"t": "La plateforme FASTR intègre une **importation directe** depuis DHIS2.", "space": 10},
        {"t": "C'est souvent la méthode la plus simple, une fois identifiés les indicateurs à inclure.", "space": 10},
        {"t": "Pour chaque pays, les données ont déjà été extraites et importées avant l'atelier.", "space": 0},
     ]},

    {"type": "image_text", "title": "Évaluation, ajustement et analyse de la qualité",
     "image": "s22_Picture_1.png", "side": "right", "tw": 4.6, "blocks": [
        {"t": "Un **outil web** pour évaluer, ajuster et analyser la qualité des données de routine.", "space": 10},
        {"t": "Importer et analyser des données de sources variées, dont DHIS2, avec des méthodes statistiques intégrées.", "space": 10},
        {"t": "Une interface conviviale, et des options souples pour visualiser et exporter les résultats.", "space": 0},
     ]},

    {"type": "content", "title": "L'assistant IA", "tw": 10.8, "blocks": [
        {"t": "La plateforme comprend un **assistant IA** qui aide, à la demande, à interpréter les données et à générer des rapports.", "size": 16, "space": 12},
        {"t": "De nombreux systèmes de santé ont plus de données que de capacités pour les analyser.", "size": 15, "bullet": True},
        {"t": "Le personnel de S&E dispose souvent de peu de temps pour des analyses approfondies.", "size": 15, "bullet": True},
        {"t": "Les compétences analytiques varient selon les équipes et les régions.", "size": 15, "bullet": True},
        {"t": "Transformer des données en récits demande des connaissances techniques et contextuelles.", "size": 15, "bullet": True},
    ]},

    {"type": "ai_flow", "title": "Comment fonctionne l'assistant IA",
     "principle": "L'IA suit le principe **« lire avant de répondre »** — elle ne devine jamais.",
     "rows": [
        {"label": "Question sur les données", "color": NAVY,
         "steps": ["Recherche la métrique", "Lit les valeurs réelles", "Répond par une visualisation"]},
        {"label": "Question de méthode", "color": DEEP_GREEN,
         "steps": ["Recherche la documentation", "Lit les détails", "Explique en langage clair"]},
     ]},

    {"type": "contrast", "title": "Un accélérateur, pas un décideur", "cols": [
        {"head": "Vous gardez le contrôle", "accent": DEEP_GREEN, "items": [
            "**Jugement** — décider ce qui est important",
            "**Interprétation** — comprendre le contexte",
            "**Action** — prendre les décisions"]},
        {"head": "L'IA interprète et explique", "accent": GREEN, "items": [
            "Les chiffres viennent de **méthodes validées** : détection des valeurs aberrantes, estimations de couverture, scores de qualité",
            "Ces calculs utilisent des **formules statistiques éprouvées**, pas l'IA",
            "À vous de décider et d'agir"]},
    ], "bottom": [
        ("Assistant générique", "Interroge le web ; pas nécessairement ancré dans les données et les méthodes d'un projet", INK3),
        ("Assistant FASTR", "Ancré dans les données et les méthodes du projet, il vérifie la source avant de répondre", GREEN),
    ]},

    {"type": "demo", "eyebrow": "Démonstration en direct  ·  ~5 min", "title": "La plateforme FASTR, en direct",
     "lead": "Passons à la plateforme pour voir le flux de bout en bout :",
     "cues": [
        "Importer les données depuis DHIS2 et parcourir les indicateurs",
        "Lancer une analyse et explorer les visualisations (qualité, volumes, couverture)",
        "Interroger l'assistant IA sur un résultat, puis générer un extrait de rapport",
     ]},

    {"type": "closing", "title": "L'approche FASTR",
     "tagline": "Analyser · Apprendre · Renforcer · Agir"},
]


def main():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    for c in CONTENT:
        BUILDERS[c["type"]](prs, c)
    total = len(CONTENT)
    no_footer = {"cover", "section", "statement", "closing"}
    for n, (c, slide) in enumerate(zip(CONTENT, prs.slides), 1):
        _fade(slide)
        if c["type"] not in no_footer:
            _footer(slide, n, total)
    out = str(REPO / "decks" / "fr" / "approche_fastr.pptx")
    prs.save(out)
    print(f"Wrote {len(CONTENT)} slides -> {out}")


if __name__ == "__main__":
    main()
