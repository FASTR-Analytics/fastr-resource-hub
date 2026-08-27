#!/usr/bin/env python3
"""Build a *natively editable* PPTX for the FASTR advocacy deck.

Marp's experimental ``--pptx-editable`` round-trips through LibreOffice and
produces files PowerPoint flags for repair. This builder skips that entirely:
it lays out real text boxes and shapes with python-pptx, so every slide opens
clean and is fully editable (text, colours, positions).

It mirrors the content and design of ``decks/<lang>/plaidoyer_plateforme_fastr.md``
(theme: fastr) — the markdown remains the source for the PDF; this script is the
source for the editable PPTX.

Usage:
    python3 decks/build_deck_pptx.py --lang fr
    python3 decks/build_deck_pptx.py --lang en
"""
from __future__ import annotations

import argparse
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parent.parent
RES = REPO / "resources"
FONT = "Poppins"

# ---- FASTR palette -------------------------------------------------------
DEEP_GREEN = RGBColor(0x09, 0x54, 0x4F)
DARK_GREEN = RGBColor(0x0C, 0x71, 0x6B)
GREEN = RGBColor(0x1F, 0x9A, 0x9C)
LIME = RGBColor(0xD0, 0xCB, 0x17)
NAVY = RGBColor(0x21, 0x56, 0x8C)
GOLD = RGBColor(0xD8, 0xA8, 0x22)
PURPLE = RGBColor(0x7A, 0x1F, 0x6E)
INK = RGBColor(0x1A, 0x1F, 0x1E)
INK2 = RGBColor(0x5A, 0x65, 0x62)
INK3 = RGBColor(0x97, 0xA0, 0x9D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_IN = 914400
SW, SH = 13.333, 7.5
LEFT = 0.62
CW = SW - 2 * LEFT

PILL_COLORS = {"navy": NAVY, "deep_green": DEEP_GREEN, "green": GREEN, "gold": GOLD}
RC_COLORS = {"navy": NAVY, "deep_green": DEEP_GREEN, "green": GREEN, "gold": GOLD}


# ---- low-level helpers ---------------------------------------------------
def _runs(text: str):
    """Split **bold** spans into (segment, bold) runs."""
    out, i = [], 0
    for m in re.finditer(r"\*\*(.+?)\*\*", text):
        if m.start() > i:
            out.append((text[i:m.start()], False))
        out.append((m.group(1), True))
        i = m.end()
    if i < len(text):
        out.append((text[i:], False))
    return out or [(text, False)]


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tb, tf


def _para(tf, text, *, size, color, first=False, bold=False, italic=False,
          align=PP_ALIGN.LEFT, bullet=False, bullet_color=LIME, space=8,
          line=1.06):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    p.space_before = Pt(0)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if bullet:
        r = p.add_run()
        r.text = "•  "
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.bold = True
        r.font.color.rgb = bullet_color
    for seg, b in _runs(text):
        r = p.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.bold = bool(b or bold)
        r.font.italic = italic
        r.font.color.rgb = color
    return p


def _rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    return sp


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _fullbleed(slide, img):
    slide.shapes.add_picture(str(img), 0, 0, Inches(SW), Inches(SH))


def _fit(img_path, max_w, max_h):
    with Image.open(img_path) as im:
        w, h = im.size
    r = w / h
    R = max_w / max_h
    if r > R:
        return max_w, max_w / r
    return max_h * r, max_h


def _title(slide, text, *, color=DEEP_GREEN, rule=GREEN, y=0.55):
    # ~54 chars fit on one line at 30pt Poppins bold across CW; titles in the
    # 48-56 char zone are ambiguous — keep deck titles out of that range.
    lines = max(1, -(-len(text) // 54))  # ceil
    th = 0.62 * lines
    _tb(slide, LEFT, y, CW, th + 0.15)[1] and None
    tb, tf = _tb(slide, LEFT, y, CW, th + 0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _para(tf, text, size=30, color=color, first=True, bold=True, space=0, line=1.02)
    rule_y = y + th + 0.06
    rule_w = min(CW, 0.62 + 0.145 * len(text)) if lines == 1 else min(CW, 8.5)
    _rect(slide, LEFT, rule_y, rule_w, 0.045, fill=rule)
    return rule_y + 0.28  # body start y


def _new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---- slide archetypes ----------------------------------------------------
def s_cover(prs, c):
    s = _new(prs)
    _fullbleed(s, RES / "backgrounds" / "cover_slide_clean.png")
    s.shapes.add_picture(str(RES / "logos" / "GFF_Logo_trimmed.png"),
                         Inches(0.8), Inches(0.55), height=Inches(0.42))
    tb, tf = _tb(s, 0.8, 2.7, 10.5, 2.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, c["title"], size=40, color=WHITE, first=True, bold=True, space=10, line=1.05)
    _para(tf, c["subtitle"], size=18, color=RGBColor(0xEC, 0xF0, 0xEF), space=0, line=1.2)
    y = SH - 1.05
    s.shapes.add_picture(str(RES / "logos" / "FASTR_White_Horiz.png"),
                         Inches(0.8), Inches(y), height=Inches(0.5))
    s.shapes.add_picture(str(RES / "logos" / "usefuldata600w.png"),
                         Inches(3.0), Inches(y + 0.07), height=Inches(0.34))


def s_section(prs, c):
    s = _new(prs)
    _fullbleed(s, RES / "backgrounds" / "section_slide.png")
    tb, tf = _tb(s, 1.4, 0, SW - 2.8, SH)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, c["title"], size=34, color=WHITE, first=True, bold=True,
          align=PP_ALIGN.CENTER, space=0, line=1.08)


def s_content(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    tb, tf = _tb(s, LEFT, by, CW, SH - by - 0.7)
    first = True
    for b in c["blocks"]:
        bullet = b.get("bullet", False)
        _para(tf, b["t"], size=16, color=INK, first=first, bullet=bullet,
              space=10 if not bullet else 6, line=1.12)
        first = False


def s_centered(prs, c):
    s = _new(prs)
    tb, tf = _tb(s, 1.1, 0, SW - 2.2, SH)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, c["title"], size=30, color=DEEP_GREEN, first=True, bold=True,
          align=PP_ALIGN.CENTER, space=18, line=1.08)
    first = False
    for p in c["paras"]:
        _para(tf, p, size=17, color=INK, align=PP_ALIGN.CENTER, space=10, line=1.2)
    # centered rule under title
    _rect(s, (SW - 3.2) / 2, 2.62, 3.2, 0.045, fill=GREEN)


def s_two_col(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    colw = (CW - 0.7) / 2
    for i, col in enumerate(c["cols"]):
        x = LEFT + i * (colw + 0.7)
        tb, tf = _tb(s, x, by, colw, SH - by - 1.1)
        _para(tf, col["head"], size=16, color=DEEP_GREEN, first=True, bold=True,
              space=8, line=1.05)
        for it in col["items"]:
            _para(tf, it, size=15, color=INK, bullet=True, space=6, line=1.12)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_columns3(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    gap = 0.5
    colw = (CW - 2 * gap) / 3
    for i, card in enumerate(c["cards"]):
        x = LEFT + i * (colw + gap)
        # accent bar left of the card
        _rect(s, x, by + 0.02, 0.06, 2.4, fill=card.get("accent", DEEP_GREEN))
        tb, tf = _tb(s, x + 0.22, by, colw - 0.22, 3.4)
        _para(tf, card["eyebrow"].upper(), size=10.5, color=card.get("accent", DEEP_GREEN),
              first=True, bold=True, space=6)
        _para(tf, card["head"], size=17, color=INK, bold=True, space=8, line=1.05)
        _para(tf, card["text"], size=13.5, color=INK2, space=8, line=1.18)
        if card.get("pill"):
            label, color = card["pill"]
            _pill(s, x + 0.22, by + 2.55, label, PILL_COLORS[color])
    if c.get("footer"):
        _callout(s, c["footer"])


def s_split(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    # right green panel, bleeding to the right/bottom edges
    px = 8.5
    _rect(s, px, by - 0.1, SW - px, SH - by + 0.1, fill=DEEP_GREEN)
    tb, tf = _tb(s, LEFT, by + 0.1, px - LEFT - 0.5, SH - by - 0.8)
    first = True
    for p in c["left"]:
        _para(tf, p, size=15.5, color=INK, first=first, space=12, line=1.2)
        first = False
    tb2, tf2 = _tb(s, px + 0.45, by, SW - px - 0.8, SH - by)
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for label, sub in c["panel"]:
        _para(tf2, label, size=17, color=LIME, first=first, bold=True, space=2, line=1.05)
        _para(tf2, sub, size=13, color=WHITE, space=14, line=1.1)
        first = False


def s_output(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    img = RES / "screenshots" / c["image"]
    max_w, max_h = 7.5, SH - by - 0.6
    w, h = _fit(img, max_w, max_h)
    s.shapes.add_picture(str(img), Inches(LEFT), Inches(by + (max_h - h) / 2),
                         Inches(w), Inches(h))
    tx = LEFT + max_w + 0.45
    tb, tf = _tb(s, tx, by, SW - tx - 0.55, SH - by - 0.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for p in c["paras"]:
        _para(tf, p, size=14.5, color=INK, first=first, space=12, line=1.22)
        first = False


def s_architecture(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    tb, tf = _tb(s, LEFT, by, 6.7, SH - by - 1.0)
    first = True
    for p in c["left"]:
        _para(tf, p, size=15.5, color=INK, first=first, space=12, line=1.2)
        first = False
    # native instance/projects diagram on the right
    ix, iy, iw, ih = 7.7, by + 0.1, 5.0, 3.7
    outer = _rect(s, ix, iy, iw, ih, fill=RGBColor(0xE8, 0xF4, 0xF3),
                  line=DEEP_GREEN, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lb, lf = _tb(s, ix + 0.25, iy + 0.16, iw - 0.5, 0.4)
    _para(lf, c["instance_label"], size=13, color=DEEP_GREEN, first=True, bold=True)
    pw = (iw - 0.5 - 2 * 0.25) / 3
    for i in range(3):
        x = ix + 0.25 + i * (pw + 0.25)
        _rect(s, x, iy + 0.8, pw, ih - 1.15, fill=WHITE, line=GREEN, line_w=1,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        pb, pf = _tb(s, x, iy + 0.8, pw, ih - 1.15)
        pf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _para(pf, c["project_label"], size=11.5, color=INK2, first=True, bold=True,
              align=PP_ALIGN.CENTER, line=1.05)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_results_chain(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    n = len(c["steps"])
    arrow = 0.42
    boxw = (CW - (n - 1) * arrow) / n
    boxy, boxh = by + 0.35, 3.15
    for i, st in enumerate(c["steps"]):
        x = LEFT + i * (boxw + arrow)
        col = RC_COLORS[st["color"]]
        _rect(s, x, boxy, boxw, boxh, fill=RGBColor(0xFA, 0xFA, 0xFA))
        _rect(s, x, boxy, boxw, 0.08, fill=col)
        tb, tf = _tb(s, x + 0.18, boxy + 0.24, boxw - 0.36, boxh - 0.48)
        _para(tf, st["eyebrow"].upper(), size=11, color=col, first=True, bold=True,
              align=PP_ALIGN.CENTER, space=6)
        _para(tf, st["title"], size=15, color=INK, bold=True, align=PP_ALIGN.CENTER,
              space=6, line=1.05)
        _para(tf, st["desc"], size=11, color=INK2, italic=True, align=PP_ALIGN.CENTER,
              space=10, line=1.12)
        _para(tf, st["pill"], size=10, color=col, bold=True, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ab, af = _tb(s, x + boxw, boxy, arrow, boxh)
            af.vertical_anchor = MSO_ANCHOR.MIDDLE
            _para(af, "→", size=22, color=INK3, first=True, align=PP_ALIGN.CENTER)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_closing(prs, c):
    s = _new(prs)
    _bg(s, DEEP_GREEN)
    tb, tf = _tb(s, 0.9, 0.75, 11, 1.0)
    _para(tf, c["title"], size=40, color=WHITE, first=True, bold=True, space=0)
    _rect(s, 0.95, 1.72, 1.6, 0.05, fill=LIME)
    tb2, tf2 = _tb(s, 0.9, 2.0, 11.2, 1.4)
    _para(tf2, c["tagline"], size=22, color=WHITE, first=True, bold=True, space=0, line=1.12)
    _rect(s, 0.95, 3.35, 8.0, 0.045, fill=LIME)
    y = 4.0
    s.shapes.add_picture(str(RES / "logos" / "FASTR_White_Horiz.png"),
                         Inches(0.95), Inches(y), height=Inches(0.5))
    s.shapes.add_picture(str(RES / "logos" / "GFF_Logo_trimmed.png"),
                         Inches(2.9), Inches(y + 0.1), height=Inches(0.3))


def _pill(slide, x, y, text, color):
    w = 0.35 + 0.11 * len(text)
    sp = _rect(slide, x, y, w, 0.32, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.word_wrap = False
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    _para(tf, text, size=10, color=WHITE, first=True, bold=True, align=PP_ALIGN.CENTER, space=0)


def _callout(slide, text):
    y = SH - 0.95
    _rect(slide, LEFT, y, CW, 0.62, fill=RGBColor(0xEA, 0xF3, 0xF3))
    _rect(slide, LEFT, y, 0.07, 0.62, fill=DEEP_GREEN)
    tb, tf = _tb(slide, LEFT + 0.25, y, CW - 0.4, 0.62)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, text, size=12, color=DEEP_GREEN, first=True, italic=True, bold=True,
          space=0, line=1.1)


def s_image(prs, c):
    """Title + optional lead sentence + one large centered image
    (path relative to resources/)."""
    s = _new(prs)
    by = _title(s, c["title"])
    if c.get("lead"):
        tb, tf = _tb(s, LEFT, by, CW, 0.75)
        _para(tf, c["lead"], size=14.5, color=INK, first=True, space=0, line=1.2)
        by += 0.85
    img = RES / c["image"]
    pad = 0.75 if c.get("footer") else 0.45
    max_w, max_h = CW, SH - by - pad
    w, h = _fit(img, max_w, max_h)
    s.shapes.add_picture(str(img), Inches(LEFT + (max_w - w) / 2),
                         Inches(by + (max_h - h) / 2), Inches(w), Inches(h))
    if c.get("footer"):
        _callout(s, c["footer"])


BUILDERS = {
    "cover": s_cover, "section": s_section, "content": s_content,
    "centered": s_centered, "two_col": s_two_col, "columns3": s_columns3,
    "split": s_split, "output": s_output, "architecture": s_architecture,
    "results_chain": s_results_chain, "closing": s_closing, "image": s_image,
}


def build(content, out_path):
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    for c in content:
        BUILDERS[c["type"]](prs, c)
    prs.save(out_path)
    return len(content)


# ---- content (mirrors the markdown deck) ---------------------------------
FR = [
    {"type": "cover", "title": "La plateforme d'analyse FASTR",
     "subtitle": "Un même endroit pour rassembler, analyser et partager les données de santé d'un pays"},
    {"type": "content", "title": "Le point de départ", "blocks": [
        {"t": "Les données de santé existent déjà. Chaque mois, les établissements les saisissent dans DHIS2. Les enquêtes auprès des formations sanitaires en produisent d'autres. Les enquêtes ménages aussi."},
        {"t": "Le problème n'est pas le manque de données. C'est qu'elles restent **dispersées, difficiles à recouper et longues à transformer en analyses**."},
        {"t": "Entre le moment où une donnée est saisie et le moment où elle éclaire une décision, il se passe souvent des semaines de travail manuel."},
    ]},
    {"type": "centered", "title": "Et si les données d'un pays vivaient — et travaillaient — au même endroit ?",
     "paras": ["Rassemblées, actualisées, prêtes à être analysées et partagées, sans repartir de zéro à chaque fois."]},
    {"type": "section", "title": "Qu'est-ce que FASTR ?"},
    {"type": "content", "title": "Une plateforme en ligne, deux rôles", "blocks": [
        {"t": "FASTR est un outil en ligne. Il ne demande aucune installation : un navigateur suffit, et l'interface existe en **français, anglais et portugais**."},
        {"t": "Il réunit deux choses que l'on trouve d'habitude dans des outils séparés :"},
        {"t": "**Un dépôt central** où vivent les données de santé d'un pays", "bullet": True},
        {"t": "**Un moteur d'analyse** qui les traite automatiquement, sans écrire de code", "bullet": True},
        {"t": "C'est cette combinaison qui distingue FASTR d'un simple entrepôt de données."},
    ]},
    {"type": "architecture", "title": "Une architecture simple : l'instance et les projets",
     "instance_label": "Instance — le pays", "project_label": "Projet",
     "left": [
         "L'**instance** est l'espace du pays. Elle contient une seule fois la structure sanitaire, les définitions d'indicateurs et toutes les sources de données. C'est la source de vérité partagée.",
         "Les **projets** sont des espaces d'analyse ciblés. Chacun prend un extrait de l'instance — une période, des régions, des indicateurs — pour répondre à une question précise.",
     ],
     "footer": "Une base commune, plusieurs analyses. Tout le monde part des mêmes données."},
    {"type": "section", "title": "Elle se connecte à vos sources de données"},
    {"type": "columns3", "title": "Trois types de sources, une même plateforme", "cards": [
        {"eyebrow": "Données de routine", "head": "SNIS / DHIS2", "accent": NAVY,
         "text": "Les statistiques mensuelles des établissements : consultations, vaccinations, accouchements.",
         "pill": ("DHIS2", "navy")},
        {"eyebrow": "Établissements", "head": "Enquêtes FOSA", "accent": DEEP_GREEN,
         "text": "L'évaluation des formations sanitaires : disponibilité des services, équipements, personnel.",
         "pill": ("HFA", "deep_green")},
        {"eyebrow": "Équité", "head": "Enquêtes ménages", "accent": GOLD,
         "text": "Les estimations de couverture par quintile de richesse, issues des enquêtes DHS et MICS.",
         "pill": ("ICEH", "gold")},
    ], "footer": "Des sources qui se complètent : le routinier, le structurel et l'équité, côte à côte."},
    {"type": "content", "title": "La connexion à DHIS2, sans ressaisie", "blocks": [
        {"t": "FASTR se connecte directement à DHIS2. On sélectionne les indicateurs et la période, et les données arrivent dans la plateforme."},
        {"t": "Cet import peut se faire de trois façons :"},
        {"t": "**Immédiatement**, à la demande", "bullet": True},
        {"t": "**À une date planifiée**, pour un chargement unique", "bullet": True},
        {"t": "**De façon récurrente**, chaque semaine ou toutes les deux semaines", "bullet": True},
        {"t": "Une fois la connexion enregistrée, les mises à jour ne demandent plus de manipulation. Les données restent alignées sur DHIS2."},
    ]},
    {"type": "split", "title": "Un dépôt central, une seule version des faits", "left": [
        "Les données d'un pays cessent d'être éparpillées entre des fichiers, des postes et des versions différentes.",
        "Elles sont importées une fois, au niveau de l'instance, et deviennent disponibles pour toutes les analyses.",
        "Quand une correction est faite à la source, elle profite à tout le monde, pas à une copie isolée.",
    ], "panel": [("Une base", "partagée par toutes les équipes"),
                 ("Un historique", "suivi dans le temps"),
                 ("Une référence", "sur laquelle chacun s'appuie")]},
    {"type": "section", "title": "Pas seulement stocker — analyser"},
    {"type": "output", "title": "Un moteur d'analyse intégré", "image": "platform_fr/modules_fr.png", "paras": [
        "Stocker des données ne répond à aucune question. FASTR va plus loin : il **analyse**.",
        "Des **modules** traitent les données automatiquement — qualité, ajustement, utilisation des services, couverture. Chacun exécute des méthodes éprouvées et produit des résultats prêts à visualiser.",
        "L'utilisateur **n'écrit aucune ligne de code**. Les résultats se calculent dans un paquet de résultats ; il le rattache et travaille.",
    ]},
    {"type": "two_col", "title": "Des méthodes reconnues, intégrées à l'outil", "cols": [
        {"head": "Ce que les modules produisent", "items": [
            "Évaluation de la **qualité des données**",
            "**Ajustement** des données incomplètes",
            "Analyse de l'**utilisation des services** et des ruptures",
            "Estimation de la **couverture** et des dénominateurs"]},
        {"head": "Ce qui les rend fiables", "items": [
            "Des méthodes **standardisées**, identiques d'un pays à l'autre",
            "Des versions **suivies**, pour savoir quel calcul a produit quel résultat",
            "Un **recalcul automatique** dès que les données changent"]},
    ], "footer": "La méthode ne dépend plus de la personne qui la connaît. Elle est dans l'outil."},
    {"type": "centered", "title": "Ce qui est puissant : rien ne se fige", "paras": [
        "Les analyses restent liées aux données, pas figées dans un fichier.",
        "Quand les données sont mises à jour, les résultats se recalculent et les graphiques suivent. Pas de copier-coller à refaire, pas de chiffre oublié.",
    ]},
    {"type": "section", "title": "Voir les résultats"},
    {"type": "two_col", "title": "Des visualisations adaptées à chaque question", "cols": [
        {"head": "Quatre formes", "items": [
            "**Graphiques** pour comparer entre catégories",
            "**Séries temporelles** pour suivre une évolution",
            "**Cartes** pour voir les écarts entre régions",
            "**Tableaux** pour le détail chiffré exact"]},
        {"head": "Que l'on maîtrise", "items": [
            "**Filtrer** et **désagréger** par région, type d'établissement, période",
            "**Personnaliser** l'apparence",
            "**Exporter** en image ou en données pour un usage externe"]},
    ], "footer": "La forme suit la question : « comment se comparent nos régions ? » n'appelle pas le même graphique que « quel est le chiffre exact ? »."},
    {"type": "output", "title": "Une analyse, pas seulement un graphique",
     "image": "platform_fr/example-viz-timeseries-fr.png", "paras": [
        "Chaque visualisation s'appuie sur un module. Ici, l'analyse des **perturbations de services** compare le volume observé au volume attendu, pour une dizaine d'indicateurs à la fois.",
        "Les zones colorées signalent les écarts. C'est un signal à recouper avec le terrain, pas une conclusion en soi.",
     ]},
    {"type": "section", "title": "Partager, à chaque public son format"},
    {"type": "columns3", "title": "Trois façons de partager les résultats", "cards": [
        {"eyebrow": "En direct", "head": "Tableaux de bord", "accent": DEEP_GREEN,
         "text": "Plusieurs visualisations sur une page, toujours à jour. Publiables via un lien public : les partenaires les ouvrent dans un navigateur, sans compte FASTR."},
        {"eyebrow": "En réunion", "head": "Présentations", "accent": GREEN,
         "text": "Des diaporamas assemblés dans la plateforme, avec pages de titre et sections. Export vers PowerPoint ou PDF pour une présentation en personne."},
        {"eyebrow": "À l'écrit", "head": "Rapports", "accent": GOLD,
         "text": "Des documents narratifs mêlant texte et chiffres en direct. Export vers Word ou PDF pour une lecture complète."},
    ], "footer": "Un même jeu de données alimente les trois. On ne refait pas le travail à chaque format."},
    {"type": "output", "title": "Le fil rouge : des chiffres toujours à jour",
     "image": "platform_fr/creating-a-report-fr.png", "paras": [
        "Dans un tableau de bord ou un rapport, les graphiques ne sont pas des captures collées qui vieillissent.",
        "Ce sont des **figures vivantes**, reliées aux données du projet. Quand les données changent, le document reflète la nouvelle réalité.",
     ]},
    {"type": "content", "title": "Un assistant IA pour interpréter", "blocks": [
        {"t": "Un assistant intégré aide à lire et interpréter les résultats. Il comprend les modules, les indicateurs et les visualisations du projet."},
        {"t": "On lui pose des questions en langage courant — « que montre la tendance de la CPN1 ? », « quelles régions ont la couverture la plus basse ? » — et il répond à partir des **données réelles du projet**, pas de suppositions."},
        {"t": "Il aide aussi à rédiger le texte des rapports, à partir d'une bibliothèque de prompts que l'équipe peut enrichir et partager."},
    ]},
    {"type": "section", "title": "Un outil pour les équipes"},
    {"type": "split", "title": "Conçu pour travailler ensemble", "left": [
        "Le travail s'organise en projets, et chaque personne reçoit un rôle adapté.",
        "Les droits se règlent finement : consulter, éditer, administrer. Un projet terminé peut être **verrouillé** pour préserver son état tout en restant consultable.",
        "Chacun voit la même donnée, dans la langue de son choix.",
    ], "panel": [("Consultation", "lire et exporter"),
                 ("Édition", "créer visualisations et rapports"),
                 ("Administration", "paramètres et accès")]},
    {"type": "results_chain", "title": "De bout en bout, un seul flux", "steps": [
        {"eyebrow": "Importer", "title": "Rassembler", "desc": "Les sources arrivent dans l'instance", "pill": "DHIS2 · HFA · ICEH", "color": "navy"},
        {"eyebrow": "Analyser", "title": "Traiter", "desc": "Les modules calculent automatiquement", "pill": "Qualité · Couverture", "color": "deep_green"},
        {"eyebrow": "Visualiser", "title": "Voir", "desc": "Graphiques, cartes, tableaux", "pill": "Explorer", "color": "green"},
        {"eyebrow": "Partager", "title": "Diffuser", "desc": "Tableaux de bord, présentations, rapports", "pill": "Décider", "color": "gold"},
    ], "footer": "Chaque étape s'enchaîne dans le même outil. La donnée ne quitte jamais la plateforme jusqu'à la décision."},
    {"type": "content", "title": "Ce que cela change", "blocks": [
        {"t": "Sans FASTR, la donnée est dispersée, retravaillée à la main et vite périmée."},
        {"t": "Avec FASTR, elle est **rassemblée** en un lieu, **analysée** avec des méthodes reconnues, et **partagée** dans le format qui convient à chaque public."},
        {"t": "Le pays passe d'une collection de fichiers à une ressource commune, à jour et fiable, au service de la décision."},
    ]},
    {"type": "closing", "title": "FASTR",
     "tagline": "Des données de routine aux décisions, sans quitter la plateforme"},
]

EN = [
    {"type": "cover", "title": "The FASTR analytics platform",
     "subtitle": "One place to bring together, analyze, and share a country's health data"},
    {"type": "content", "title": "Where we start", "blocks": [
        {"t": "Health data already exists. Every month, facilities enter it into DHIS2. Facility assessments produce more. Household surveys, more still."},
        {"t": "The problem isn't a shortage of data. It's that it stays **scattered, hard to reconcile, and slow to turn into analysis**."},
        {"t": "Between the moment data is entered and the moment it informs a decision, weeks of manual work often pass."},
    ]},
    {"type": "centered", "title": "What if a country's data lived — and worked — in one place?",
     "paras": ["Brought together, kept current, ready to analyze and share, without starting from scratch every time."]},
    {"type": "section", "title": "What is FASTR?"},
    {"type": "content", "title": "An online platform, two roles", "blocks": [
        {"t": "FASTR is an online tool. Nothing to install: a browser is enough, and the interface comes in **French, English, and Portuguese**."},
        {"t": "It brings together two things usually found in separate tools:"},
        {"t": "**A central repository** where a country's health data lives", "bullet": True},
        {"t": "**An analysis engine** that processes it automatically, with no code to write", "bullet": True},
        {"t": "It's this combination that sets FASTR apart from a plain data warehouse."},
    ]},
    {"type": "architecture", "title": "A simple architecture: the instance and projects",
     "instance_label": "Instance — the country", "project_label": "Project",
     "left": [
         "The **instance** is the country's space. It holds the health structure, the indicator definitions, and every data source — once. It is the shared source of truth.",
         "**Projects** are focused analysis spaces. Each takes a slice of the instance — a period, some regions, some indicators — to answer one specific question.",
     ],
     "footer": "One shared base, many analyses. Everyone starts from the same data."},
    {"type": "section", "title": "It connects to your data sources"},
    {"type": "columns3", "title": "Three kinds of source, one platform", "cards": [
        {"eyebrow": "Routine data", "head": "HMIS / DHIS2", "accent": NAVY,
         "text": "Facilities' monthly statistics: consultations, vaccinations, deliveries.",
         "pill": ("DHIS2", "navy")},
        {"eyebrow": "Facilities", "head": "Facility surveys (HFA)", "accent": DEEP_GREEN,
         "text": "The health facility assessment: service availability, equipment, staffing.",
         "pill": ("HFA", "deep_green")},
        {"eyebrow": "Equity", "head": "Household surveys", "accent": GOLD,
         "text": "Coverage estimates by wealth quintile, drawn from DHS and MICS surveys.",
         "pill": ("ICEH", "gold")},
    ], "footer": "Sources that complement each other: the routine, the structural, and the equity view, side by side."},
    {"type": "content", "title": "Connecting to DHIS2, without re-keying", "blocks": [
        {"t": "FASTR connects directly to DHIS2. You pick the indicators and the period, and the data flows into the platform."},
        {"t": "That import can run three ways:"},
        {"t": "**Immediately**, on demand", "bullet": True},
        {"t": "**At a scheduled time**, for a one-off load", "bullet": True},
        {"t": "**On a recurring schedule**, weekly or every two weeks", "bullet": True},
        {"t": "Once the connection is saved, updates take no further handling. The data stays aligned with DHIS2."},
    ]},
    {"type": "split", "title": "A central repository, a single version of the facts", "left": [
        "A country's data stops being scattered across files, machines, and differing versions.",
        "It is imported once, at the instance level, and becomes available to every analysis.",
        "When a correction is made at the source, everyone benefits — not one isolated copy.",
    ], "panel": [("One base", "shared by every team"),
                 ("One history", "tracked over time"),
                 ("One reference", "everyone relies on")]},
    {"type": "section", "title": "Not just storage — analysis"},
    {"type": "output", "title": "A built-in analysis engine", "image": "platform_en/modules_en.png", "paras": [
        "Storing data answers no question. FASTR goes further: it **analyzes**.",
        "**Modules** process the data automatically — quality, adjustment, service utilization, coverage. Each runs proven methods and produces results ready to visualize.",
        "The user **writes no code**. The results compute in a results package; they attach it and work.",
    ]},
    {"type": "two_col", "title": "Recognized methods, built into the tool", "cols": [
        {"head": "What the modules produce", "items": [
            "**Data quality** assessment",
            "**Adjustment** of incomplete data",
            "**Service utilization** and disruption analysis",
            "**Coverage** and denominator estimation"]},
        {"head": "What makes them reliable", "items": [
            "**Standardized** methods, the same from one country to the next",
            "**Tracked** versions, so you know which computation produced which result",
            "**Automatic recomputation** whenever the data changes"]},
    ], "footer": "The method no longer depends on the person who knows it. It's in the tool."},
    {"type": "centered", "title": "The powerful part: nothing gets frozen", "paras": [
        "Analyses stay tied to the data, not frozen in a file.",
        "When the data is updated, the results recompute and the charts follow. No copy-paste to redo, no forgotten figure.",
    ]},
    {"type": "section", "title": "Seeing the results"},
    {"type": "two_col", "title": "Visualizations suited to each question", "cols": [
        {"head": "Four forms", "items": [
            "**Charts** to compare across categories",
            "**Time series** to follow a trend",
            "**Maps** to see regional gaps",
            "**Tables** for the exact figures"]},
        {"head": "That you control", "items": [
            "**Filter** and **disaggregate** by region, facility type, period",
            "**Customize** the appearance",
            "**Export** as an image or data for outside use"]},
    ], "footer": "Form follows the question: “how do our regions compare?” doesn't call for the same chart as “what's the exact number?”"},
    {"type": "output", "title": "An analysis, not just a chart",
     "image": "platform_en/example-viz-timeseries-en.png", "paras": [
        "Every visualization builds on a module. Here, the **service disruption** analysis compares observed volume to expected volume, across a dozen indicators at once.",
        "The shaded areas flag the gaps. That's a signal to triangulate with the field, not a conclusion on its own.",
     ]},
    {"type": "section", "title": "Sharing — the right format for each audience"},
    {"type": "columns3", "title": "Three ways to share the results", "cards": [
        {"eyebrow": "Live", "head": "Dashboards", "accent": DEEP_GREEN,
         "text": "Several visualizations on one page, always current. Publishable via a public link: partners open them in a browser, with no FASTR account."},
        {"eyebrow": "In the room", "head": "Presentations", "accent": GREEN,
         "text": "Slide decks assembled in the platform, with title and section pages. Export to PowerPoint or PDF for an in-person talk."},
        {"eyebrow": "In writing", "head": "Reports", "accent": GOLD,
         "text": "Narrative documents blending prose and live figures. Export to Word or PDF for a full read."},
    ], "footer": "One dataset feeds all three. You don't redo the work for each format."},
    {"type": "output", "title": "The through-line: figures always current",
     "image": "platform_en/creating-a-report-en.png", "paras": [
        "In a FASTR dashboard or report, the charts aren't pasted-in pictures that go stale.",
        "They are **live figures**, tied to the project's data. When the data changes, the document reflects the new reality.",
     ]},
    {"type": "content", "title": "An AI assistant to interpret", "blocks": [
        {"t": "A built-in assistant helps read and interpret the results. It understands the project's modules, indicators, and visualizations."},
        {"t": "You ask questions in plain language — “what does the ANC1 trend show?”, “which regions have the lowest coverage?” — and it answers from the **project's real data**, not guesses."},
        {"t": "It also helps draft report text, from a prompt library the team can extend and share."},
    ]},
    {"type": "section", "title": "A tool for teams"},
    {"type": "split", "title": "Built to work together", "left": [
        "Work is organized into projects, and each person gets a fitting role.",
        "Permissions are fine-grained: view, edit, administer. A finished project can be **locked** to preserve its state while staying viewable.",
        "Everyone sees the same data, in the language of their choice.",
    ], "panel": [("View", "read and export"),
                 ("Edit", "create visualizations and reports"),
                 ("Administer", "settings and access")]},
    {"type": "results_chain", "title": "End to end, a single flow", "steps": [
        {"eyebrow": "Import", "title": "Bring together", "desc": "Sources arrive in the instance", "pill": "DHIS2 · HFA · ICEH", "color": "navy"},
        {"eyebrow": "Analyze", "title": "Process", "desc": "Modules compute automatically", "pill": "Quality · Coverage", "color": "deep_green"},
        {"eyebrow": "Visualize", "title": "See", "desc": "Charts, maps, tables", "pill": "Explore", "color": "green"},
        {"eyebrow": "Share", "title": "Deliver", "desc": "Dashboards, presentations, reports", "pill": "Decide", "color": "gold"},
    ], "footer": "Each step follows in the same tool. The data never leaves the platform on its way to a decision."},
    {"type": "content", "title": "What this changes", "blocks": [
        {"t": "Without FASTR, data is scattered, reworked by hand, and quickly out of date."},
        {"t": "With FASTR, it is **brought together** in one place, **analyzed** with recognized methods, and **shared** in the format each audience needs."},
        {"t": "A country moves from a collection of files to a shared resource — current and trustworthy — in the service of decisions."},
    ]},
    {"type": "closing", "title": "FASTR",
     "tagline": "From routine data to decisions, without leaving the platform"},
]

CONTENT = {"fr": FR, "en": EN}

# ---- governance deck (mirrors decks/<lang>/platform_security_sustainability.md)
GOV_EN = [
    {"type": "cover", "title": "Security, costs, and ownership",
     "subtitle": "How the FASTR platform protects a country's data — and what it takes to run it"},
    {"type": "section", "title": "Data security"},
    {"type": "content", "title": "What the platform holds", "blocks": [
        {"t": "**Monthly service totals per facility** — for example, 45 first antenatal visits at one clinic in March", "bullet": True},
        {"t": "The **same figures as the DHIS2 reports** the ministry already produces", "bullet": True},
        {"t": "**No patient records** — no names, no addresses, nothing individual", "bullet": True},
        {"t": "This is **today's scope, not a ceiling**: hosting individual-level records later is not excluded — it would come with a further, defined set of safeguards (technical annex)", "bullet": True},
    ]},
    {"type": "image", "title": "What \"secure\" means for a platform like this",
     "lead": "In software and data hosting, security is not one thing — it is five. The next slides take them one at a time: the risk, then what answers it.",
     "image": "diagrams/gov_security_dimensions.png"},
    {"type": "image", "title": "Could our data leave the country's control?",
     "lead": "**Sovereignty.** The risk: data flowing into a shared pool, or visible to another country. The answer: there is no shared pool — each country runs its own installation, sharing is technically impossible, and everything can be exported in full at any time.",
     "image": "diagrams/gov_country_isolation.png"},
    {"type": "content", "title": "Could someone gain unauthorized access?", "blocks": [
        {"t": "**Confidentiality.** The risk: a leaked password, or an insider reaching beyond their role. The answer: accounts, roles, and identity re-checked at every single request."},
        {"t": "The ministry decides **who gets an account**, and each person's role", "bullet": True},
        {"t": "Roles set what a person can do — **view, edit, or administer** — and in which projects", "bullet": True},
        {"t": "Identity is **re-checked on every request**, not just at login", "bullet": True},
        {"t": "Only a **small, named group of administrators** can change the setup", "bullet": True},
    ]},
    {"type": "image", "title": "Could data be intercepted — or lost?",
     "lead": "**Confidentiality and availability.** The risk: interception on the network, or a server failure taking the data with it. The answer: nothing readable leaves the platform, and automatic copies mean everything can be rebuilt.",
     "image": "diagrams/gov_security_layers.png"},
    {"type": "content", "title": "Could the figures be altered?", "blocks": [
        {"t": "**Integrity.** The risk: figures altered without record, or analyses that cannot be reproduced. The answer: one source of truth, recorded imports, versioned results."},
        {"t": "**DHIS2 remains the source of truth** — re-imported months are refreshed to match it", "bullet": True},
        {"t": "**Every import is recorded**: an indicator-by-indicator ledger shows months loaded, when, and failures", "bullet": True},
        {"t": "**Results come in versioned, dated packages** — the same package gives the same numbers to everyone", "bullet": True},
        {"t": "Analysis methods and parameters are **documented and logged with each package**", "bullet": True},
    ]},
    {"type": "image", "title": "Could the AI expose what it sees?",
     "lead": "**Confidentiality and accountability.** The risk: an assistant that sees too much. The answer: it is handed aggregated totals only — within that user's own permissions, and every question is logged. **Never the underlying rows.**",
     "image": "diagrams/gov_ai_boundary.png"},
    {"type": "section", "title": "Costs and ownership"},
    {"type": "image", "title": "Running costs",
     "lead": "Three lines: **server hosting** (fixed), **AI use** (metered), **maintenance** (shared team). No license fees, no per-user fees — the software is open source.",
     "image": "diagrams/gov_cost_structure.png"},
    {"type": "content", "title": "The figures — planning estimates", "blocks": [
        {"t": "Estimates at 2026 prices and today's usage. **Not a quote.**"},
        {"t": "**Hosting** a country's instance: about **USD 500–800 per year** — larger countries toward the top", "bullet": True},
        {"t": "**AI use**: about **USD 350 per month** (~USD 4,000 per year) for a typical country; the largest countries several times this", "bullet": True},
        {"t": "**Shared platform maintenance**: roughly **USD 7,000–10,000 per country per year** today — falls as more countries join", "bullet": True},
        {"t": "**Optional support**: about **USD 5,000–15,000 per year**, depending on level", "bullet": True},
        {"t": "**Total: about USD 12,000 per year** — up to USD 20,000–30,000 with support", "bullet": True},
        {"t": "Lines include the external services: the **login service (Clerk)** and the **AI API (Anthropic)** — no hidden subscriptions", "bullet": True},
        {"t": "**Economies of scale are already working**: one team, one codebase, every country — the maintenance share falls as countries join, economics a country gives up by hosting alone", "bullet": True},
    ]},
    {"type": "content", "title": "Hosting — and what moving really takes", "blocks": [
        {"t": "Hosted **centrally today**, while the platform is in active development — every country receives fixes and new features the same day", "bullet": True},
        {"t": "Built **portable**: the same software runs unchanged on a ministry's own servers — no rebuild needed", "bullet": True},
        {"t": "**But migration is a project, not a switch**: readiness assessment, server procurement, team training, a parallel run, then cutover — measured in **months, planned jointly**", "bullet": True},
        {"t": "**Hosting alone means carrying alone** what is shared today: maintenance, monitoring, upgrades and fixes become the ministry team's own workload", "bullet": True},
        {"t": "**All of a country's data can be exported in full at any time**, regardless of hosting", "bullet": True},
    ]},
    {"type": "image", "title": "The path to country ownership",
     "lead": "For countries who do want to host themselves, there will be effort required from the Ministry of Health. A team will need to run servers, backups, and upgrades; and procurement in place for hosting and for AI services, which are bought from a commercial provider (e.g. Anthropic, OpenAI).",
     "image": "diagrams/gov_ownership_roadmap.png"},
    {"type": "content", "title": "The essentials", "blocks": [
        {"t": "**Aggregated totals today** — no patient records; any future individual-level hosting comes with added safeguards", "bullet": True},
        {"t": "**One installation per country** — no pathway between countries", "bullet": True},
        {"t": "**About USD 12,000 per year** to run — hosting, metered AI, shared maintenance; no license or per-user fees", "bullet": True},
        {"t": "**Country hosting by 2030** — the stated transition goal", "bullet": True},
    ]},
    {"type": "section", "title": "Technical annex — for IT and DHIS2 teams"},
    {"type": "content", "title": "Architecture and isolation", "blocks": [
        {"t": "Each country instance is a **separate Docker deployment**: its own application container, **PostgreSQL** database, **Valkey** cache, and storage volume, on a private network", "bullet": True},
        {"t": "Within an instance, each project has its **own database** (project_{uuid}) for authoring content", "bullet": True},
        {"t": "Computed results live in **versioned results packages** (parquet, queried via DuckDB) — projects read only the package attached to them", "bullet": True},
    ]},
    {"type": "content", "title": "Authentication and access control", "blocks": [
        {"t": "Identity managed by **Clerk**; session claims **verified by middleware on every request** — auth bypass disabled in production", "bullet": True},
        {"t": "**Two-tier RBAC**: instance-level permissions plus per-project roles resolved from the database on each request", "bullet": True},
        {"t": "Server access: **SSH by cryptographic key only**, restricted to two named engineers", "bullet": True},
    ]},
    {"type": "content", "title": "Data protection and operations", "blocks": [
        {"t": "**TLS terminates at nginx**, certificates auto-renewed via certbot; secrets injected at runtime — never in images, never in the browser", "bullet": True},
        {"t": "**R analysis modules run in ephemeral containers** that mount only their own run's sandbox", "bullet": True},
        {"t": "Backups: **pg_dump every 30 minutes** (3-day window) plus **volume snapshots daily/weekly/monthly**; restores permission-gated with a server-held key", "bullet": True},
    ]},
    {"type": "content", "title": "AI integration, precisely", "blocks": [
        {"t": "Claude (Anthropic) reached through a **server-side proxy**; the API key never leaves the server", "bullet": True},
        {"t": "Tool calls run **in the user's authenticated session** — the AI holds no permissions of its own", "bullet": True},
        {"t": "Data tools return **aggregated metric outputs only**: no facility-identifier dimension; dimensions with more than 20 values summarized as counts", "bullet": True},
        {"t": "**Server-side web search/fetch (Anthropic-hosted) is enabled** in the project chat; every request logged (user, project, model, tokens) with daily and weekly limits", "bullet": True},
    ]},
    {"type": "content", "title": "Data and interoperability", "blocks": [
        {"t": "Source data: **aggregated DHIS2 numerators** (facility-month counts), imported server-side with a per-indicator import ledger", "bullet": True},
        {"t": "**Full export at any time**, regardless of hosting arrangement", "bullet": True},
        {"t": "The codebase is **open source**: github.com/FASTR-Analytics/platform", "bullet": True},
    ]},
    {"type": "content", "title": "If patient-level data were hosted: the requirements", "blocks": [
        {"t": "Aggregated-only is today's design choice. Individual-level hosting is not excluded — and would be conditional on a further set of safeguards:"},
        {"t": "**Legal basis first**: national data-protection and health-data law, data-sharing agreements — possibly **in-country hosting** as a precondition", "bullet": True},
        {"t": "**Encryption at rest** in addition to transit; **pseudonymization** wherever analysis allows", "bullet": True},
        {"t": "**Minimum-necessary access**: finer roles, per-field restrictions, audit trails on every record access", "bullet": True},
        {"t": "Formal **incident-response and breach-notification** procedures, independent security testing, certification-level practices", "bullet": True},
        {"t": "**Ministry governance**: a designated data-access authority deciding who may see what, and why", "bullet": True},
    ]},
]

GOV_FR = [
    {"type": "cover", "title": "Sécurité, coûts et appropriation",
     "subtitle": "Comment la plateforme FASTR protège les données d'un pays — et ce qu'il faut pour la faire tourner"},
    {"type": "section", "title": "Sécurité des données"},
    {"type": "content", "title": "Ce que contient la plateforme", "blocks": [
        {"t": "Des **totaux mensuels de services par établissement** — par exemple, 45 premières consultations prénatales dans une clinique en mars", "bullet": True},
        {"t": "Les **mêmes chiffres que les rapports DHIS2** que le ministère produit déjà", "bullet": True},
        {"t": "**Aucun dossier patient** — pas de noms, pas d'adresses, rien d'individuel", "bullet": True},
        {"t": "C'est **le périmètre d'aujourd'hui, pas un plafond** : héberger plus tard des données individuelles n'est pas exclu — cela viendrait avec un ensemble supplémentaire de garanties, définies (annexe technique)", "bullet": True},
    ]},
    {"type": "image", "title": "Ce que « sécurisé » veut dire pour une telle plateforme",
     "lead": "En logiciel et en hébergement de données, la sécurité n'est pas une chose — c'en est cinq. Les diapositives suivantes les prennent une à une : le risque, puis ce qui y répond.",
     "image": "diagrams_fr/gov_security_dimensions.png"},
    {"type": "image", "title": "Nos données peuvent-elles échapper au contrôle du pays ?",
     "lead": "**Souveraineté.** Le risque : des données versées dans un espace partagé, ou visibles d'un autre pays. La réponse : il n'existe aucun espace partagé — chaque pays a sa propre installation, le partage est techniquement impossible, et tout peut être exporté intégralement à tout moment.",
     "image": "diagrams_fr/gov_country_isolation.png"},
    {"type": "content", "title": "Un accès non autorisé est-il possible ?", "blocks": [
        {"t": "**Confidentialité.** Le risque : un mot de passe qui fuite, ou un initié qui dépasse son rôle. La réponse : des comptes, des rôles, et une identité revérifiée à chaque requête."},
        {"t": "Le ministère décide **qui reçoit un compte**, et le rôle de chaque personne", "bullet": True},
        {"t": "Le rôle fixe ce qu'une personne peut faire — **consulter, éditer ou administrer** — et dans quels projets", "bullet": True},
        {"t": "L'identité est **revérifiée à chaque requête**, pas seulement à la connexion", "bullet": True},
        {"t": "Seul un **petit groupe d'administrateurs identifiés** peut modifier la configuration", "bullet": True},
    ]},
    {"type": "image", "title": "Les données peuvent-elles être interceptées — ou perdues ?",
     "lead": "**Confidentialité et disponibilité.** Le risque : une interception sur le réseau, ou une panne de serveur qui emporte les données. La réponse : rien de lisible ne quitte la plateforme, et des copies automatiques permettent de tout reconstruire.",
     "image": "diagrams_fr/gov_security_layers.png"},
    {"type": "content", "title": "Les chiffres peuvent-ils être altérés ?", "blocks": [
        {"t": "**Intégrité.** Le risque : des chiffres modifiés sans trace, ou des analyses impossibles à reproduire. La réponse : une source de vérité, des importations enregistrées, des résultats versionnés."},
        {"t": "**DHIS2 reste la source de vérité** — les mois réimportés sont rafraîchis pour y correspondre", "bullet": True},
        {"t": "**Chaque importation est enregistrée** : un registre par indicateur montre les mois chargés, quand, et les échecs", "bullet": True},
        {"t": "**Les résultats arrivent en paquets versionnés et datés** — le même paquet donne les mêmes chiffres à tous", "bullet": True},
        {"t": "Méthodes et paramètres d'analyse **documentés et consignés avec chaque paquet**", "bullet": True},
    ]},
    {"type": "image", "title": "L'IA peut-elle exposer ce qu'elle voit ?",
     "lead": "**Confidentialité et traçabilité.** Le risque : un assistant qui en voit trop. La réponse : des totaux agrégés uniquement — dans la limite des permissions de l'utilisateur, chaque question journalisée. **Jamais les lignes en dessous.**",
     "image": "diagrams_fr/gov_ai_boundary.png"},
    {"type": "section", "title": "Coûts et appropriation"},
    {"type": "image", "title": "Coûts de fonctionnement",
     "lead": "Trois lignes : **hébergement du serveur** (fixe), **usage de l'IA** (au compteur), **maintenance** (équipe partagée). Ni licence ni frais par utilisateur — le logiciel est open source.",
     "image": "diagrams_fr/gov_cost_structure.png"},
    {"type": "content", "title": "Les chiffres — estimations de planification", "blocks": [
        {"t": "Estimations aux prix 2026 et à l'usage actuel. **Ce n'est pas un devis.**"},
        {"t": "**Hébergement** de l'instance : environ **500 à 800 USD par an** — les grands pays vers le haut", "bullet": True},
        {"t": "**Usage de l'IA** : environ **350 USD par mois** (~4 000 USD par an) pour un pays typique ; plusieurs fois plus pour les plus grands", "bullet": True},
        {"t": "**Maintenance partagée** : environ **7 000 à 10 000 USD par pays et par an** aujourd'hui — en baisse à mesure que des pays rejoignent", "bullet": True},
        {"t": "**Appui optionnel** : environ **5 000 à 15 000 USD par an** selon le niveau", "bullet": True},
        {"t": "**Total : environ 12 000 USD par an** — jusqu'à 20 000–30 000 USD avec l'appui", "bullet": True},
        {"t": "Les lignes incluent les services externes : le **service de connexion (Clerk)** et l'**API d'IA (Anthropic)** — aucun abonnement caché", "bullet": True},
        {"t": "**Les économies d'échelle jouent déjà** : une équipe, un code, tous les pays — la part de maintenance baisse à mesure que des pays rejoignent, des économies qu'un pays abandonne en s'hébergeant seul", "bullet": True},
    ]},
    {"type": "content", "title": "Hébergement — et ce que migrer demande vraiment", "blocks": [
        {"t": "Hébergée **de façon centralisée aujourd'hui**, pendant le développement actif — chaque pays reçoit correctifs et nouveautés le jour même", "bullet": True},
        {"t": "Construite **portable** : le même logiciel tourne tel quel sur les serveurs d'un ministère — aucune reconstruction", "bullet": True},
        {"t": "**Mais migrer est un projet, pas un interrupteur** : évaluation de préparation, achat des serveurs, formation, fonctionnement en parallèle, puis bascule — en **mois, planifiés ensemble**", "bullet": True},
        {"t": "**S'héberger seul, c'est porter seul** ce qui est partagé aujourd'hui : maintenance, supervision, mises à niveau et correctifs reviennent à l'équipe du ministère", "bullet": True},
        {"t": "**Toutes les données d'un pays peuvent être exportées intégralement à tout moment**, quel que soit l'hébergement", "bullet": True},
    ]},
    {"type": "image", "title": "Le chemin vers l'appropriation par le pays",
     "lead": "Pour les pays qui souhaitent s'héberger eux-mêmes, un effort du ministère de la Santé sera nécessaire. Une équipe devra faire tourner serveurs, sauvegardes et mises à niveau ; et des marchés devront être en place pour l'hébergement et pour les services d'IA, achetés auprès d'un fournisseur commercial (par ex. Anthropic, OpenAI).",
     "image": "diagrams_fr/gov_ownership_roadmap.png"},
    {"type": "content", "title": "L'essentiel", "blocks": [
        {"t": "**Des totaux agrégés aujourd'hui** — aucun dossier patient ; tout hébergement futur de données individuelles s'accompagnerait de garanties supplémentaires", "bullet": True},
        {"t": "**Une installation par pays** — aucun passage entre pays", "bullet": True},
        {"t": "**Environ 12 000 USD par an** de fonctionnement — hébergement, IA au compteur, maintenance partagée ; ni licence ni frais par utilisateur", "bullet": True},
        {"t": "**Hébergement par le pays d'ici 2030** — l'objectif de transition affiché", "bullet": True},
    ]},
    {"type": "section", "title": "Annexe technique — équipes informatiques et DHIS2"},
    {"type": "content", "title": "Architecture et isolation", "blocks": [
        {"t": "Chaque instance pays est un **déploiement Docker séparé** : conteneur applicatif, base **PostgreSQL**, cache **Valkey** et volume de stockage propres, sur un réseau privé", "bullet": True},
        {"t": "Au sein d'une instance, chaque projet a sa **propre base** (project_{uuid}) pour le contenu d'édition", "bullet": True},
        {"t": "Les résultats calculés résident dans des **paquets de résultats versionnés** (parquet, interrogés via DuckDB) — chaque projet ne lit que le paquet qui lui est rattaché", "bullet": True},
    ]},
    {"type": "content", "title": "Authentification et contrôle d'accès", "blocks": [
        {"t": "Identité gérée par **Clerk** ; jetons de session **vérifiés par middleware à chaque requête** — contournement d'authentification désactivé en production", "bullet": True},
        {"t": "**RBAC à deux niveaux** : permissions d'instance et rôles par projet, résolus en base à chaque requête", "bullet": True},
        {"t": "Accès aux serveurs : **SSH par clé cryptographique uniquement**, restreint à deux ingénieurs nommés", "bullet": True},
    ]},
    {"type": "content", "title": "Protection des données et exploitation", "blocks": [
        {"t": "**TLS terminé sur nginx**, certificats renouvelés via certbot ; secrets injectés à l'exécution — jamais dans les images, jamais côté navigateur", "bullet": True},
        {"t": "**Modules R exécutés dans des conteneurs éphémères** qui ne montent que le répertoire de leur exécution", "bullet": True},
        {"t": "Sauvegardes : **pg_dump toutes les 30 minutes** (fenêtre de 3 jours) plus **instantanés de volume quotidiens/hebdomadaires/mensuels** ; restauration contrôlée par permission et clé côté serveur", "bullet": True},
    ]},
    {"type": "content", "title": "L'intégration IA, précisément", "blocks": [
        {"t": "Claude (Anthropic) atteint via un **proxy côté serveur** ; la clé d'API ne quitte jamais le serveur", "bullet": True},
        {"t": "Les appels d'outils s'exécutent **dans la session authentifiée de l'utilisateur** — l'IA n'a aucune permission propre", "bullet": True},
        {"t": "Outils de données : **sorties métriques agrégées uniquement** — aucune dimension identifiant d'établissement ; toute dimension de plus de 20 valeurs résumée en nombre", "bullet": True},
        {"t": "**Recherche web côté serveur (Anthropic) activée** dans le chat projet ; chaque requête journalisée (utilisateur, projet, modèle, jetons) avec limites quotidiennes et hebdomadaires", "bullet": True},
    ]},
    {"type": "content", "title": "Données et interopérabilité", "blocks": [
        {"t": "Données source : **numérateurs DHIS2 agrégés** (totaux établissement-mois), importés côté serveur avec registre d'importation par indicateur", "bullet": True},
        {"t": "**Export intégral à tout moment**, quel que soit l'hébergement", "bullet": True},
        {"t": "Code **open source** : github.com/FASTR-Analytics/platform", "bullet": True},
    ]},
    {"type": "content", "title": "Si des données patients étaient hébergées : les exigences", "blocks": [
        {"t": "Le périmètre agrégé est un choix de conception, valable aujourd'hui. Héberger des données individuelles n'est pas exclu — et serait conditionné à des garanties supplémentaires :"},
        {"t": "**La base légale d'abord** : droit national des données de santé, accords de partage — possiblement **l'hébergement dans le pays** comme précondition", "bullet": True},
        {"t": "**Chiffrement au repos** en plus du transit ; **pseudonymisation** partout où l'analyse le permet", "bullet": True},
        {"t": "**Accès au strict nécessaire** : rôles plus fins, restrictions par champ, journaux d'audit sur chaque accès", "bullet": True},
        {"t": "Procédures formelles de **réponse aux incidents et de notification des violations**, tests de sécurité indépendants, pratiques de niveau certification", "bullet": True},
        {"t": "**Gouvernance ministérielle** : une autorité désignée d'accès aux données décidant qui voit quoi, et pourquoi", "bullet": True},
    ]},
]

DECKS = {
    "plaidoyer": {"content": CONTENT, "file": "plaidoyer_plateforme_fastr.pptx"},
    "governance": {"content": {"en": GOV_EN, "fr": GOV_FR},
                   "file": "platform_security_sustainability.pptx"},
}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--lang", default="fr", choices=("en", "fr"))
    ap.add_argument("--deck", default="plaidoyer", choices=sorted(DECKS))
    ap.add_argument("--out")
    args = ap.parse_args()
    deck = DECKS[args.deck]
    out = args.out or str(REPO / "decks" / args.lang / deck["file"])
    n = build(deck["content"][args.lang], out)
    print(f"Wrote {n} slides -> {out}")


if __name__ == "__main__":
    main()
