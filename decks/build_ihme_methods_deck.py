#!/usr/bin/env python3
"""Build a natively editable PPTX for the FASTR ↔ IHME methods-exchange deck.

Same system as build_approche_fastr.py: real editable text boxes and shapes on
the FASTR design system — brand palette, flat accent-bar layouts, section
dividers, page-numbered footers, fade transitions.

CONTENT MASTER: "/Users/claireboulange/Desktop/WORK TODAY/IHME/IHME_deck_text_MASTER.md"
The flow is one-directional: edit the master, fold changes into CONTENT below,
re-run. This script must NEVER export or overwrite the master (or any text
file) — that is how edits were lost once already.

Usage:  python3 decks/build_ihme_methods_deck.py
Deps:   pip install python-pptx
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

try:
    from pptx.enum.line import MSO_LINE
except ImportError:  # older python-pptx
    MSO_LINE = None

REPO = Path(__file__).resolve().parent.parent
RES = REPO / "resources"
FONT = "Poppins"

DEEP_GREEN = RGBColor(0x09, 0x54, 0x4F)
GREEN = RGBColor(0x1F, 0x9A, 0x9C)
LIME = RGBColor(0xD0, 0xCB, 0x17)
NAVY = RGBColor(0x21, 0x56, 0x8C)
GOLD = RGBColor(0xD8, 0xA8, 0x22)
INK = RGBColor(0x1A, 0x1F, 0x1E)
INK2 = RGBColor(0x5A, 0x65, 0x62)
INK3 = RGBColor(0x97, 0xA0, 0x9D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_IN = 914400
SW, SH = 13.333, 7.5
LEFT = 0.62
CW = SW - 2 * LEFT
STEP_COLORS = [NAVY, DEEP_GREEN, GREEN, GOLD]
TINT = {NAVY: RGBColor(0xED, 0xF1, 0xF7), DEEP_GREEN: RGBColor(0xE8, 0xF1, 0xF0),
        GREEN: RGBColor(0xE9, 0xF4, 0xF4), GOLD: RGBColor(0xF9, 0xF3, 0xE4)}
CENTER = PP_ALIGN.CENTER
MID = MSO_ANCHOR.MIDDLE


# ---- helpers -------------------------------------------------------------
def _runs(text):
    out, i = [], 0
    for m in re.finditer(r"\*\*(.+?)\*\*", text):
        if m.start() > i:
            out.append((text[i:m.start()], False))
        out.append((m.group(1), True))
        i = m.end()
    if i < len(text):
        out.append((text[i:], False))
    return out or [(text, False)]


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tb, tf


def _para(tf, text, *, size, color, first=False, bold=False, italic=False,
          align=PP_ALIGN.LEFT, bullet=False, bullet_color=LIME, space=8, line=1.1):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    p.space_before = Pt(0)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if bullet:
        r = p.add_run()
        r.text = "•  "
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.bold = True
        r.font.color.rgb = bullet_color
    for seg, b in _runs(text):
        r = p.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.bold = bool(b or bold)
        r.font.italic = italic
        r.font.color.rgb = color
    return p


def _rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    el = sp._element
    st = el.find(qn("p:style"))
    if st is not None:
        el.remove(st)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    return sp


def _seg(slide, x1, y1, x2, y2, color, w=2.0, dash=False):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.shadow.inherit = False
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    if dash and MSO_LINE is not None:
        try:
            ln.line.dash_style = MSO_LINE.DASH
        except Exception:
            pass
    return ln


def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _fullbleed(slide, img):
    slide.shapes.add_picture(str(img), 0, 0, Inches(SW), Inches(SH))


def _blocks(tf, blocks, base=14):
    first = True
    for b in blocks:
        bl = b.get("bullet", False)
        _para(tf, b["t"], size=b.get("size", base), color=b.get("color", INK), first=first,
              bullet=bl, bold=b.get("bold", False), italic=b.get("italic", False),
              space=b.get("space", 10 if not bl else 6), line=b.get("line", 1.2))
        first = False


def _title(slide, text, *, y=0.5):
    lines = max(1, -(-len(text) // 50))
    th = 0.6 * lines
    _, tf = _tb(slide, LEFT, y, CW, th + 0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _para(tf, text, size=30, color=DEEP_GREEN, first=True, bold=True, space=0, line=1.02)
    ry = y + th + 0.06
    rw = min(CW, 0.6 + 0.14 * len(text)) if lines == 1 else min(CW, 8.5)
    _rect(slide, LEFT, ry, rw, 0.045, fill=GREEN)
    return ry + 0.3


def _callout(slide, text):
    y = SH - 1.02
    _rect(slide, LEFT, y, 0.06, 0.5, fill=DEEP_GREEN)
    _, tf = _tb(slide, LEFT + 0.24, y, CW - 0.35, 0.5)
    tf.vertical_anchor = MID
    _para(tf, text, size=12, color=DEEP_GREEN, first=True, italic=True, bold=True, space=0, line=1.15)


def _footer(slide, n, total):
    _rect(slide, LEFT, SH - 0.44, CW, 0.012, fill=RGBColor(0xE4, 0xE7, 0xE5))
    _, lf = _tb(slide, LEFT, SH - 0.38, CW - 1.2, 0.3)
    _para(lf, "FASTR × IHME  ·  METHODS EXCHANGE", size=8, color=INK3, first=True, bold=True, space=0)
    _, pf = _tb(slide, SW - LEFT - 1.0, SH - 0.38, 1.0, 0.3)
    _para(pf, f"{n} / {total}", size=8, color=INK3, first=True, bold=True, align=PP_ALIGN.RIGHT, space=0)


def _fade(slide, spd="med"):
    sld = slide._element
    for t in sld.findall(qn("p:transition")):
        sld.remove(t)
    sld.append(parse_xml(
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'spd="{spd}"><p:fade/></p:transition>'))


def _new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---- archetypes ----------------------------------------------------------
def s_cover(prs, c):
    s = _new(prs)
    _fullbleed(s, RES / "backgrounds" / "cover_slide_clean.png")
    s.shapes.add_picture(str(RES / "logos" / "GFF_Logo_trimmed.png"), Inches(0.8), Inches(0.55), height=Inches(0.42))
    _, tf = _tb(s, 0.8, 2.3, 11.6, 2.9)
    tf.vertical_anchor = MID
    _para(tf, c["title"], size=33, color=WHITE, first=True, bold=True, space=12, line=1.06)
    _para(tf, c["subtitle"], size=17, color=RGBColor(0xEC, 0xF0, 0xEF), space=8, line=1.25)
    _para(tf, c["meta"], size=13, color=RGBColor(0xC9, 0xD6, 0xD3), space=0, line=1.2)
    y = SH - 1.05
    s.shapes.add_picture(str(RES / "logos" / "FASTR_White_Horiz.png"), Inches(0.8), Inches(y), height=Inches(0.5))
    s.shapes.add_picture(str(RES / "logos" / "usefuldata600w.png"), Inches(3.0), Inches(y + 0.07), height=Inches(0.34))


def s_section(prs, c):
    s = _new(prs)
    _fullbleed(s, RES / "backgrounds" / "section_slide.png")
    _, tf = _tb(s, 1.2, 0, SW - 2.4, SH)
    tf.vertical_anchor = MID
    _para(tf, c["title"], size=32, color=WHITE, first=True, bold=True, align=CENTER,
          space=10 if c.get("kicker") else 0, line=1.1)
    if c.get("kicker"):
        _para(tf, c["kicker"], size=14, color=LIME, align=CENTER, space=0, line=1.25)


def s_blocks_cards(prs, c):
    """Text paragraphs with one or two accent cards along the bottom."""
    s = _new(prs)
    by = _title(s, c["title"])
    cards = c["cards"]
    card_h = c.get("card_h", 1.1)
    _, tf = _tb(s, LEFT, by, c.get("tw", CW), SH - by - card_h - 0.85)
    tf.vertical_anchor = MID
    _blocks(tf, c["blocks"], base=14)
    cy = SH - card_h - 0.62
    gap = 0.5
    cw2 = (CW - gap * (len(cards) - 1)) / len(cards)
    for i, (tag, desc, acc) in enumerate(cards):
        x = LEFT + i * (cw2 + gap)
        _rect(s, x, cy, cw2, card_h, fill=None, line=acc, line_w=1.25)
        _rect(s, x, cy, 0.06, card_h, fill=acc)
        _, cf = _tb(s, x + 0.28, cy, cw2 - 0.45, card_h)
        cf.vertical_anchor = MID
        _para(cf, tag, size=13, color=acc, first=True, bold=True, space=3)
        _para(cf, desc, size=11.5, color=INK, space=0, line=1.14)


def s_cycle(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    n = len(c["nodes"])
    arrow = 0.5
    bw = (CW - (n - 1) * arrow) / n
    bh, y0 = 3.2, by + 0.4
    for i, (label, desc) in enumerate(c["nodes"]):
        x = LEFT + i * (bw + arrow)
        col = STEP_COLORS[i % 4]
        _rect(s, x, y0, bw, bh, fill=TINT[col])
        _rect(s, x, y0, bw, 0.09, fill=col)
        _, tf = _tb(s, x + 0.24, y0 + 0.34, bw - 0.48, bh - 0.5)
        _para(tf, label, size=17, color=col, first=True, bold=True, align=CENTER, space=10, line=1.05)
        _para(tf, desc, size=11.5, color=INK2, align=CENTER, space=0, line=1.24)
        if i < n - 1:
            _, af = _tb(s, x + bw, y0, arrow, bh)
            af.vertical_anchor = MID
            _para(af, "→", size=22, color=INK3, first=True, align=CENTER)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_quad(prs, c):
    """2×2 grid of accent blocks, with an optional framing line above the grid."""
    s = _new(prs)
    by = _title(s, c["title"])
    gap = 0.3
    cw2 = (CW - gap) / 2
    default_rh = 2.18
    if c.get("intro"):
        ilines = max(1, -(-len(c["intro"]) // 150))
        ih = 0.24 * ilines + 0.14
        _, itf = _tb(s, LEFT, by, CW, ih)
        _para(itf, c["intro"], size=12.5, color=INK2, first=True, italic=True, space=0, line=1.18)
        by += ih
        default_rh = 2.18 - ih / 2
    rh = c.get("rh", default_rh)
    for i, (head, desc) in enumerate(c["items"]):
        col = STEP_COLORS[i % 4]
        x = LEFT + (i % 2) * (cw2 + gap)
        y = by + 0.03 + (i // 2) * (rh + gap)
        _rect(s, x, y, cw2, rh, fill=TINT[col])
        _rect(s, x, y, 0.07, rh, fill=col)
        _, tf = _tb(s, x + 0.3, y + 0.2, cw2 - 0.55, rh - 0.36)
        _para(tf, head, size=14, color=col, first=True, bold=True, space=6, line=1.05)
        _para(tf, desc, size=11, color=INK, space=0, line=1.24)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_numbered(prs, c):
    """Numbered workstreams as full-width rows — survives long descriptions."""
    s = _new(prs)
    by = _title(s, c["title"])
    _, tf = _tb(s, LEFT, by + 0.05, CW, SH - by - 1.2)
    first = True
    for i, (head, desc) in enumerate(c["steps"], 1):
        col = STEP_COLORS[(i - 1) % 4]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        p.line_spacing = 1.12
        r = p.add_run(); r.text = f"0{i}   "
        r.font.size = Pt(13); r.font.name = FONT; r.font.bold = True; r.font.color.rgb = col
        r = p.add_run(); r.text = head
        r.font.size = Pt(12.5); r.font.name = FONT; r.font.bold = True; r.font.color.rgb = INK
        _para(tf, desc, size=10.5, color=INK2, space=9, line=1.18)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_table(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    if c.get("intro"):
        ilines = max(1, -(-len(c["intro"]) // 115))
        ih = 0.26 * ilines + 0.08
        _, itf = _tb(s, LEFT, by, CW, ih)
        _para(itf, c["intro"], size=13, color=INK2, first=True, italic=True, space=0, line=1.15)
        by += ih + 0.08
    rows, cols = len(c["rows"]) + 1, len(c["headers"])
    th = rows * min(0.72, (SH - by - 1.15) / rows)
    gt = s.shapes.add_table(rows, cols, Inches(LEFT), Inches(by), Inches(CW), Inches(th)).table
    if c.get("colw"):
        for j, w in enumerate(c["colw"]):
            gt.columns[j].width = Inches(w)
    for j, htxt in enumerate(c["headers"]):
        cell = gt.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = DEEP_GREEN
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.vertical_anchor = MID
        _para(cell.text_frame, htxt, size=11.5, color=WHITE, first=True, bold=True, space=0, line=1.05)
    for i, row in enumerate(c["rows"], 1):
        for j, val in enumerate(row):
            cell = gt.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF4, 0xF7, 0xF6)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MID
            bold = (j == 0 and c.get("bold_first"))
            _para(cell.text_frame, val, size=c.get("fs", 11), color=DEEP_GREEN if bold else INK,
                  first=True, bold=bold, space=0, line=1.08)
    if c.get("note"):
        _, nf = _tb(s, LEFT, SH - 1.05, CW, 0.58)
        _para(nf, c["note"], size=9.5, color=INK3, first=True, italic=True, space=0, line=1.12)


def s_two_col(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    if c.get("lead"):
        _, lf = _tb(s, LEFT, by, CW, 0.62)
        _para(lf, c["lead"], size=13.5, color=INK, first=True, space=0, line=1.2)
        by += 0.68
    widths = c.get("widths", [1, 1])
    tot = sum(widths)
    x = LEFT
    for i, col in enumerate(c["cols"]):
        w = (CW - 0.7) * widths[i] / tot
        _, tf = _tb(s, x, by, w, SH - by - 0.7)
        _para(tf, col["head"], size=14.5, color=DEEP_GREEN, first=True, bold=True, space=10, line=1.08)
        if col.get("numbered"):
            for k, it in enumerate(col["items"], 1):
                p = tf.add_paragraph()
                p.space_after = Pt(8)
                p.line_spacing = 1.2
                r = p.add_run(); r.text = f"{k}.  "
                r.font.size = Pt(col.get("size", 12)); r.font.name = FONT
                r.font.bold = True; r.font.color.rgb = GREEN
                for seg, b in _runs(it):
                    r = p.add_run(); r.text = seg
                    r.font.size = Pt(col.get("size", 12)); r.font.name = FONT
                    r.font.bold = b; r.font.color.rgb = INK
        else:
            for it in col["items"]:
                _para(tf, it, size=col.get("size", 12), color=INK, bullet=col.get("bullet", True),
                      space=8, line=1.2)
        x += w + 0.7
    if c.get("footer"):
        _callout(s, c["footer"])


def s_ladders(prs, c):
    """Adjustment slide — lead, one full-width replacement hierarchy, rule lines."""
    s = _new(prs)
    by = _title(s, c["title"])
    _, lf = _tb(s, LEFT, by, CW, 0.62)
    _para(lf, c["lead"], size=13, color=INK, first=True, space=0, line=1.2)
    y0 = by + 0.7
    lad = c["ladder"]
    lh = 0.52 + 0.27 * len(lad["steps"])
    _rect(s, LEFT, y0, CW, lh, fill=TINT[DEEP_GREEN])
    _rect(s, LEFT, y0, CW, 0.08, fill=DEEP_GREEN)
    _, tf = _tb(s, LEFT + 0.32, y0 + 0.22, CW - 0.6, lh - 0.36)
    _para(tf, lad["head"], size=13.5, color=DEEP_GREEN, first=True, bold=True, space=7, line=1.05)
    for k, step in enumerate(lad["steps"], 1):
        p = tf.add_paragraph()
        p.space_after = Pt(3)
        p.line_spacing = 1.12
        r = p.add_run(); r.text = f"{k}.  "
        r.font.size = Pt(11.5); r.font.name = FONT; r.font.bold = True; r.font.color.rgb = DEEP_GREEN
        r = p.add_run(); r.text = step
        r.font.size = Pt(11.5); r.font.name = FONT; r.font.color.rgb = INK
    _, nf = _tb(s, LEFT, y0 + lh + 0.16, CW, SH - y0 - lh - 1.25)
    _blocks(nf, [{"t": t, "size": 11.5, "bullet": True, "space": 5, "line": 1.18} for t in c["rules"]])
    if c.get("footer"):
        _callout(s, c["footer"])


def s_cascade(prs, c):
    """Back-calculation entry + demographic cascade chain, then paragraphs."""
    s = _new(prs)
    by = _title(s, c["title"])
    _, lf = _tb(s, LEFT, by, CW, 0.9)
    _blocks(lf, [{"t": c["lead"], "size": 12.5, "space": 0, "line": 1.2}])
    ey = by + 0.98
    ew = 4.1
    _rect(s, LEFT, ey, ew, 0.52, fill=None, line=GOLD, line_w=1.25)
    _rect(s, LEFT, ey, 0.06, 0.52, fill=GOLD)
    _, ef = _tb(s, LEFT + 0.24, ey, ew - 0.4, 0.52)
    ef.vertical_anchor = MID
    _para(ef, c["entry"], size=11, color=INK, first=True, space=0, line=1.12)
    _, df = _tb(s, LEFT + 0.55, ey + 0.52, 0.5, 0.36)
    _para(df, "↓", size=15, color=GOLD, first=True, bold=True, align=CENTER, space=0)
    nodes, ops = c["nodes"], c["ops"]
    n = len(nodes)
    aw = 0.82
    bw = (CW - (n - 1) * aw) / n
    ny = ey + 0.9
    nh = 0.95
    for i, label in enumerate(nodes):
        x = LEFT + i * (bw + aw)
        col = DEEP_GREEN if i == 0 else GREEN
        _rect(s, x, ny, bw, nh, fill=TINT[col])
        _rect(s, x, ny, bw, 0.07, fill=col)
        _, tf = _tb(s, x + 0.06, ny + 0.09, bw - 0.12, nh - 0.14)
        tf.vertical_anchor = MID
        _para(tf, label, size=11, color=DEEP_GREEN, first=True, bold=True, align=CENTER, space=0, line=1.08)
        if i < n - 1:
            _, af = _tb(s, x + bw + 0.02, ny + 0.06, aw - 0.04, nh - 0.1)
            af.vertical_anchor = MID
            _para(af, ops[i], size=8, color=INK2, first=True, italic=True, align=CENTER, space=1, line=1.05)
            _para(af, "→", size=14, color=INK3, align=CENTER, space=0)
    _, bf = _tb(s, LEFT, ny + nh + 0.22, CW, SH - ny - nh - 0.75)
    _blocks(bf, [{"t": t, "size": 11.5, "space": 7, "line": 1.2} for t in c["paras"]])
    if c.get("footer"):
        _callout(s, c["footer"])


def s_projection(prs, c):
    """Coverage estimation slide: text left, native schematic chart right."""
    s = _new(prs)
    by = _title(s, c["title"])
    _, tf = _tb(s, LEFT, by, 6.5, SH - by - 1.15)
    _blocks(tf, c["blocks"], base=12)
    cx, cy = 7.55, by + 0.25
    cw3, ch = 4.9, 3.5
    x0, y0 = cx + 0.42, cy + ch - 0.45
    px_w, px_h = cw3 - 0.7, ch - 0.75
    _seg(s, x0, cy + 0.1, x0, y0, INK3, w=1.2)
    _seg(s, x0, y0, cx + cw3, y0, INK3, w=1.2)
    _, yl = _tb(s, cx - 0.28, cy + 0.05, 0.75, 0.3)
    _para(yl, "Coverage", size=9, color=INK3, first=True, italic=True, space=0)
    years = ["2019", "2020", "2021", "2022", "2023", "2024", "2025"]
    n = len(years)
    xs = [x0 + 0.22 + i * (px_w - 0.35) / (n - 1) for i in range(n)]
    for i, yr in enumerate(years):
        _, xf = _tb(s, xs[i] - 0.3, y0 + 0.05, 0.6, 0.25)
        _para(xf, yr, size=8, color=INK3, first=True, align=CENTER, space=0)
    def Y(v):
        return y0 - v * px_h
    hmis = [0.42, 0.36, 0.44, 0.40, 0.47, 0.52, 0.58]
    for i in range(n - 1):
        _seg(s, xs[i], Y(hmis[i]), xs[i + 1], Y(hmis[i + 1]), GREEN, w=2.4)
    refs = {1: 0.62, 4: 0.68}
    for i, v in refs.items():
        _rect(s, xs[i] - 0.055, Y(v) - 0.055, 0.11, 0.11, fill=NAVY)
    lastx, lastv = 4, refs[4]
    proj = [lastv, lastv + (hmis[5] - hmis[4]), lastv + (hmis[6] - hmis[4])]
    for i in range(2):
        _seg(s, xs[lastx + i], Y(proj[i]), xs[lastx + i + 1], Y(proj[i + 1]), GOLD, w=2.4, dash=True)
    _seg(s, xs[lastx], Y(refs[4]) + 0.07, xs[lastx], Y(hmis[4]) - 0.02, INK3, w=1.0, dash=True)
    _, gf = _tb(s, xs[lastx] + 0.06, (Y(refs[4]) + Y(hmis[4])) / 2 - 0.11, 1.35, 0.3)
    _para(gf, "gap retained", size=8, color=INK2, first=True, italic=True, space=0)
    ly = cy + ch + 0.12
    legend = [("Reference estimates (points)", NAVY), ("HMIS-derived coverage", GREEN),
              ("Projection", GOLD)]
    lx = cx
    for label, col in legend:
        _rect(s, lx, ly + 0.07, 0.26, 0.05, fill=col)
        _, lg = _tb(s, lx + 0.34, ly - 0.03, 1.85, 0.3)
        _para(lg, label, size=8, color=INK2, first=True, space=0, line=1.0)
        lx += 1.85
    if c.get("footer"):
        _callout(s, c["footer"])


def s_contrast(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    half = (CW - 0.6) / 2
    ph, y0 = 2.95, by + 0.1
    for i, col in enumerate(c["cols"]):
        x = LEFT + i * (half + 0.6)
        acc = col["accent"]
        _rect(s, x, y0 + 0.02, 0.05, ph - 0.04, fill=acc)
        _, tf = _tb(s, x + 0.26, y0, half - 0.26, ph)
        _para(tf, col["head"], size=15, color=acc, first=True, bold=True, space=10, line=1.08)
        for it in col["items"]:
            _para(tf, it, size=11.5, color=INK, bullet=True, space=8, line=1.2)
    bottom = c.get("bottom", [])
    for i, (tag, desc, acc) in enumerate(bottom):
        cw2 = (CW - 0.5 * (len(bottom) - 1)) / len(bottom)
        cy = y0 + ph + 0.3
        x = LEFT + i * (cw2 + 0.5)
        _rect(s, x, cy, cw2, 1.0, fill=None, line=acc, line_w=1.25)
        _rect(s, x, cy, 0.06, 1.0, fill=acc)
        _, tf = _tb(s, x + 0.28, cy, cw2 - 0.45, 1.0)
        tf.vertical_anchor = MID
        _para(tf, tag, size=13, color=acc, first=True, bold=True, space=3)
        _para(tf, desc, size=11.5, color=INK, space=0, line=1.14)
    if c.get("footer"):
        _callout(s, c["footer"])


def s_threecol(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    gap = 0.42
    cw3 = (CW - 2 * gap) / 3
    for i, col in enumerate(c["cols"]):
        acc = STEP_COLORS[i % 4]
        x = LEFT + i * (cw3 + gap)
        _rect(s, x, by + 0.02, cw3, 0.06, fill=acc)
        _, tf = _tb(s, x, by + 0.2, cw3, SH - by - 0.85)
        _para(tf, col["head"], size=13.5, color=acc, first=True, bold=True, space=10, line=1.1)
        for it in col["items"]:
            _para(tf, it, size=11, color=INK, bullet=True, bullet_color=acc, space=9, line=1.22)


def s_pillars3(prs, c):
    s = _new(prs)
    by = _title(s, c["title"])
    gap = 0.45
    cw3 = (CW - 2 * gap) / 3
    ch, y0 = 3.9, by + 0.2
    for i, p in enumerate(c["pillars"]):
        x = LEFT + i * (cw3 + gap)
        col = STEP_COLORS[i % 4]
        _rect(s, x, y0, cw3, ch, fill=TINT[col])
        _rect(s, x, y0, cw3, 0.09, fill=col)
        _, tf = _tb(s, x + 0.26, y0 + 0.28, cw3 - 0.5, ch - 0.45)
        _para(tf, p["kicker"].upper(), size=10.5, color=col, first=True, bold=True, space=7)
        _para(tf, p["head"], size=14.5, color=INK, bold=True, space=9, line=1.12)
        _para(tf, p["desc"], size=11.5, color=INK2, space=0, line=1.26)


def s_discussion(prs, c):
    s = _new(prs)
    _bg(s, DEEP_GREEN)
    _, tf = _tb(s, 0.9, 0.75, 11.4, 1.0)
    _para(tf, c["title"], size=38, color=WHITE, first=True, bold=True, space=0)
    _rect(s, 0.95, 1.75, 1.6, 0.05, fill=LIME)
    _, qf = _tb(s, 0.95, 2.25, 11.2, 3.6)
    for i, q in enumerate(c["questions"]):
        _para(qf, q, size=16.5, color=WHITE, first=(i == 0), bullet=True, bullet_color=LIME,
              space=16, line=1.25)
    cy = SH - 0.95
    s.shapes.add_picture(str(RES / "logos" / "FASTR_White_Horiz.png"), Inches(0.95), Inches(cy - 0.25), height=Inches(0.5))
    s.shapes.add_picture(str(RES / "logos" / "usefuldata600w.png"), Inches(3.25), Inches(cy - 0.17), height=Inches(0.34))
    s.shapes.add_picture(str(RES / "logos" / "GFF_Logo_trimmed.png"), Inches(6.15), Inches(cy - 0.15), height=Inches(0.3))


BUILDERS = {"cover": s_cover, "section": s_section, "blocks_cards": s_blocks_cards,
            "cycle": s_cycle, "quad": s_quad, "numbered": s_numbered, "table": s_table,
            "two_col": s_two_col, "ladders": s_ladders, "cascade": s_cascade,
            "projection": s_projection, "contrast": s_contrast, "threecol": s_threecol,
            "pillars3": s_pillars3, "discussion": s_discussion}


# ---- CONTENT — folded verbatim from IHME_deck_text_MASTER.md --------------
CONTENT = [
    # Slide 1
    {"type": "cover",
     "title": "FASTR: HMIS data quality assessment, adjustment, and coverage estimation",
     "subtitle": "Methods overview and reflections on the IHME administrative coverage analysis",
     "meta": "FASTR ↔ IHME technical exchange  ·  July 2026  ·  Global Financing Facility / World Bank — FASTR Analytics  ·  Presented by Dr Ashley Sheffel"},

    # Slide 2
    {"type": "blocks_cards", "title": "Objective, data and scope",
     "blocks": [
        {"t": "FASTR is a GFF initiative supporting country-led efforts to improve the timely generation and use of data for decision-making, and it rests on four technical approaches. This presentation covers the first of them, the analysis of routine HMIS data, known as RMNCAH-N service use monitoring.", "size": 13, "space": 9},
        {"t": "That analysis has three objectives: to assess data quality at national and subnational level, to track changes in the volume of priority RMNCAH-N services, and to monitor coverage progress against country targets and benchmarks.", "size": 13, "space": 9},
        {"t": "Analyses are quarterly, conducted with Ministry of Health teams, and work from monthly facility records extracted from DHIS2. Coverage is currently estimated in 24 countries, subnationally in 20 of them, typically over about five years of data.", "size": 13, "space": 9},
        {"t": "Because these are the records as countries hold them, they differ from the annual aggregates submitted to WHO and UNICEF, which are the basis of the IHME analysis.", "size": 13, "space": 9},
        {"t": "The two analyses also put different questions to the data.", "size": 13.5, "space": 0},
     ],
     "cards": [
        ("IHME analysis", "Whether administrative data compiled globally can support monitoring of vaccine coverage at national and subnational level.", NAVY),
        ("FASTR", "What a country can read from its own routine data between surveys.", GREEN),
     ],
     "notes": "The FASTR methods documentation was shared with the IHME team ahead of the meeting, and these slides follow it. [INPUT] Number of countries, if wanted in the first bullet."},

    # Slide 3
    {"type": "cycle", "title": "Analytical pipeline", "nodes": [
        ("Assess", "Completeness, outliers and internal consistency are evaluated at facility-month level, before any aggregation."),
        ("Adjust", "Flagged values are replaced using each facility's own reporting history, and four datasets are retained: unadjusted, outlier-adjusted, completeness-adjusted, and adjusted for both."),
        ("Analyze", "Service utilization is summarized over time, and a separate model identifies periods of service disruption."),
        ("Estimate", "Target population denominators are constructed, and coverage is estimated at national and subnational level."),
     ],
     "footer": "In routine country work these outputs are used mostly for service volumes and disruption. Coverage estimation is the newest part of the pipeline, and it is the part this presentation examines.",
     "notes": "Module identifiers, if asked: data quality assessment is M1, adjustment M2, service utilization M3, disruption detection M11, denominators M5 and coverage M6."},

    # Slide 4
    {"type": "table", "title": "Measures produced",
     "intro": "The pipeline produces four families of output, shown here with the level at which each is reported.",
     "headers": ["Measure", "Examples", "Level reported"], "colw": [2.6, 6.3, 3.19], "bold_first": True, "fs": 11.5,
     "rows": [
        ["Data quality metrics", "Reporting completeness, prevalence of outliers, consistency ratios between related indicators", "Facility, district and national"],
        ["Adjusted service volumes", "ANC1 and ANC4, institutional deliveries, BCG, Penta1 to Penta3, MCV1, outpatient attendance", "Facility-month and above"],
        ["Utilization and disruption", "Change in level and trend, and periods in which service delivery departed from its expected path", "Any level of aggregation"],
        ["Coverage estimates", "ANC1, ANC4, institutional delivery, skilled birth attendance and PNC1; BCG, Penta1, Penta3, OPV1 to OPV3, MCV1, MCV2 and rotavirus", "National and DHS region"],
     ],
     "note": "Coverage needs a survey value to anchor it, so it is available for fewer indicators and at fewer levels than the other measures. Further indicators can be programmed for a country, as vitamin A and the fully immunized child composite are in Nigeria. Dropout between doses and zero-dose counts are not currently produced."},

    # Slide 5
    {"type": "quad", "title": "Administrative coverage: known sources of error",
     "intro": "Administrative coverage divides the services a health system recorded by an estimate of the people who should have received them. The difficulties below arise wherever it is computed.",
     "items": [
        ("Numerators",
         "The numerator is the count of services facilities report. It is incomplete where facilities do not report, distorted by data entry errors such as an extra digit, and missing services delivered privately or in communities. A wrong count moves both the level of coverage and its trend."),
        ("Denominators",
         "The denominator is the estimated number of people who should have received the service, usually projected from the last census. Projections drift as that census ages, and drift most at subnational level, and different programs may hold different figures for the same population. A wrong population shifts the level of coverage but leaves the trend broadly intact."),
        ("Coverage validity",
         "Dividing one by the other often gives coverage above 100%, in DHIS2 data as in the eJRF data IHME examined. Such a value proves that one of the two terms is wrong, but not which. Where DTP1 exceeds 100%, zero-dose children cannot be counted at all."),
        ("What follows from each",
         "A wrong population alone leaves the direction of change readable. A wrong count spoils level and direction alike. When both are wrong, the errors can offset and produce a figure that looks plausible and is not."),
     ],
     "footer": "IHME found 45% of countries reporting national DTP1 coverage above 100% at least once between 2016 and 2024. FASTR reaches coverage differently: it cleans the service counts, then builds the denominator from those counts and a survey rather than a projection.",
     "notes": "Detail behind the blocks, if it comes up. Outliers also arise where a cumulative total is entered as a single month's value. Survey coverage cannot exceed 100%, because it is measured directly among the children in the sample, which is part of why the two sources diverge. Note that this slide describes conventional administrative coverage, which is not what FASTR computes; it is the problem the FASTR method is built to avoid."},

    # Slide 6
    {"type": "section", "title": "The FASTR approach",
     "kicker": "Numerators, denominators, and coverage estimation"},

    # Slide 7
    {"type": "table", "title": "Numerators: data quality assessment",
     "intro": "Data quality is assessed at facility-month level, before any aggregation. Gaps and outliers that are detectable at this resolution leave no trace once the data are aggregated to annual district totals.",
     "headers": ["Component", "Question", "Rule"], "colw": [1.7, 2.5, 7.89], "bold_first": True, "fs": 10.5,
     "rows": [
        ["Completeness", "Are all expected facility-months reported?",
         "The proportion of facilities reporting in each month is calculated as a share of those expected to report. Facilities that are silent for six or more consecutive months at the start or end of their series are treated as inactive rather than as missing."],
        ["Outliers", "Are values implausibly high?",
         "A value is flagged where it lies more than ten median absolute deviations above the facility median, or where a single facility-month contributes more than 80% of that facility's trailing twelve-month volume. Both criteria apply only to counts above one hundred."],
        ["Consistency", "Do related indicators behave as expected?",
         "Ratios are checked at district level, since patients often attend different facilities within a catchment. Penta1 to Penta3 and ANC1 to ANC4 are expected at or above 0.95, BCG to institutional deliveries within ±30%, and confirmed malaria cases to treatment within ±10%. Outlier months are excluded first."],
     ],
     "note": "Each facility-month also receives a composite score across tracer indicators, passing only where it clears all three checks. The ten-deviation threshold and the hundred-event floor were inherited from an earlier implementation and are under review, since a low-volume indicator can hide implausible values below that floor."},

    # Slide 8
    {"type": "ladders", "title": "Numerators: adjustment methods",
     "lead": "Flagged outliers and missing facility-months are replaced rather than dropped, so that each facility's series is complete when it is aggregated. Every replacement is computed from valid months only, meaning months that were reported and were not themselves flagged as outliers.",
     "ladder": {"head": "The first method with sufficient data applies", "steps": [
        "A centered six-month average, taken over the three months on either side, is used where the surrounding window holds valid months",
        "A forward-looking six-month window is used at the start of a series",
        "A backward-looking six-month window is used at the end",
        "The same calendar month in the previous year is used where that month was reported and unflagged",
        "The mean of all valid months for that facility and indicator is used",
        "The value remains missing where no valid replacement exists"]},
     "rules": [
        "Each window average is computed from the valid months the window contains, and where it contains none the hierarchy moves on to the next method.",
        "Where a series needs both outlier replacement and gap filling, outliers are replaced first, and the gap filling then draws only on months that were never flagged, so a gap is never filled from another imputed value.",
        "Mortality and stillbirth indicators are not adjusted, and neither is any indicator that never reaches one hundred events in a facility-month.",
        "Four datasets are carried through the pipeline, so any downstream result can be recomputed under each adjustment scenario.",
     ],
     "footer": "Across seven countries examined at national level, adjustment raised annual totals for institutional live births by a median of 1.7%, ranging from no change where reporting was near-complete to 7.1% where it was weakest.",
     "notes": "The same hierarchy governs outlier replacement and completeness imputation, with two differences: what triggers them, and the same-calendar-month step, which applies only to outlier replacement, not to gap filling."},

    # Slide 9
    {"type": "cascade", "title": "Denominators: derivation from service volumes",
     "lead": "Standard administrative coverage divides a service count by a projected population. FASTR keeps the administrative numerator and replaces only that population, deriving it from the service data themselves through one anchor indicator. If a survey put ANC1 coverage at 80% and 8,000 first antenatal visits were recorded, the implied number of pregnancies is 10,000.",
     "entry": "The anchor indicator's HMIS count is divided by its survey coverage",
     "nodes": ["Pregnancies", "Deliveries", "Births", "Live births", "DPT-eligible infants", "Measles-eligible children"],
     "ops": ["− pregnancy loss", "+ twins", "− stillbirths", "− neonatal deaths", "− post-neonatal deaths"],
     "paras": [
        "The anchor indicator, together with the cascade it generates, is referred to as a **chain**, and more than one indicator can serve as the anchor.",
        "Because every denominator descends from the anchor's count multiplied by fixed rates, coverage for any other indicator reduces to its volume over the anchor's, scaled by the anchor's survey coverage. Penta1 from an ANC1 chain therefore compares Penta1 doses with ANC1 visits, not with an independently counted population of infants.",
        "Survey inputs are specific to both year and geography. Each survey contributes a coverage value that is held until the next survey supersedes it, and where a survey publishes regional estimates, each region's denominator is derived from its own regional value.",
        "These remain administrative estimates rather than survey estimates. The survey is not measuring coverage here; it enters once, to fix how large the denominator should be.",
     ],
     "notes": "The anchor value is taken from the highest-priority source available for that indicator, which is DHS where a DHS estimate exists. The demographic rates are configured for each country, so they reflect that country's own mortality and fertility profile rather than a single global setting. Within a country analysis, each rate is a single value applied across all years and administrative units, which is the limitation stated later. If asked about numerator over-reporting: the denominator is proportional to the anchor's count in the same year, so systematic over-reporting of the anchor partially cancels in the coverage ratio rather than compounding against an independent denominator. The extent of that cancellation has not been quantified, and it forms part of the planned methods work."},

    # Slide 10
    {"type": "quad", "title": "Denominators: chain selection",
     "intro": "Several indicators could anchor the cascade. One is chosen for the country and then used everywhere.",
     "rh": 1.95,
     "items": [
        ("The candidates",
         "ANC1, institutional delivery, skilled birth attendance and Penta1 can each anchor a chain subnationally. BCG can anchor one at national level only."),
        ("How one is chosen",
         "A data check proposes the chain whose implied population sits closest to UN WPP. The analyst can accept that proposal or choose another on programmatic grounds."),
        ("Applied everywhere",
         "The chosen chain supplies the denominator for every indicator and every region. Denominators from two different chains are never mixed."),
        ("What UN WPP does",
         "UN WPP only decides which chain is proposed. It never supplies the denominator that enters the coverage calculation."),
     ],
     "footer": "The anchor indicator gets no independent estimate: dividing its count by a denominator built from that same count simply returns the survey value.",
     "notes": "Fuller detail is in the annex. The automatic comparison runs over ANC1, institutional delivery and Penta1, taking the median ratio of each candidate's implied population to UN WPP across pregnancies, live births and DPT-eligible infants, and choosing the one closest to 1.00. BCG is excluded because its cascade is national-only; skilled birth attendance is computed but is not an automatic candidate. The manual override accepts ANC1, delivery, BCG or Penta1."},

    # Slide 11
    {"type": "projection", "title": "Coverage estimation and projection",
     "blocks": [
        {"t": "Coverage is calculated as the adjusted service count divided by the chain-derived denominator for the corresponding target population.", "size": 12, "space": 8},
        {"t": "The series is anchored to the most recent reference estimate available for the indicator, and for later years the year-on-year change in HMIS-based coverage is added to that value.", "size": 12, "space": 5},
        {"t": "projected(t) = last reference value + [ HMIS coverage(t) − HMIS coverage(reference year) ]", "size": 11, "italic": True, "color": DEEP_GREEN, "space": 8},
        {"t": "Three series are reported side by side rather than combined into one: reference estimates appear in the years for which they exist, HMIS-derived coverage appears as the full ratio series and is not truncated, and the projection runs from the most recent reference year forward.", "size": 12, "space": 8},
        {"t": "Where HMIS-derived coverage and the reference estimate disagree at the same point, the difference is retained in the output rather than reconciled, since it points to a problem in the numerator or the denominator.", "size": 12, "space": 8},
        {"t": "A reference estimate is a household survey value in the years when a survey was conducted. For immunization indicators the reference series also draws on the WHO and UNICEF estimates of national immunization coverage, which fill the years between surveys and can therefore become the anchor for the projection.", "size": 12, "space": 0},
     ],
     "footer": "The projection rests on one assumption, that a change in HMIS coverage reflects a real change in coverage. Whether it does is taken up under open questions.",
     "notes": "Two things are true at once, and it is worth keeping them apart. In constructing the denominator, the reference value is carried forward unchanged until a later one supersedes it. In the output series, reference estimates are shown in the years for which they exist, because the carried value sits inside the denominators rather than being drawn as a flat line. Be ready on source composition: in the Zambia BCG example, DHS supplies 2002, 2007, 2013 and 2018, while WHO and UNICEF estimates fill the remaining years, so the 2024 anchor for the projection is a WHO and UNICEF estimate rather than a survey. The denominators in that country are unaffected, because the selected chain is institutional delivery and its reference values are DHS throughout."},

    # Slide 12
    {"type": "two_col", "title": "Subnational estimation: approach and requirements",
     "lead": "Subnational estimation uses the same method as national estimation. Only the inputs change: each region's own service counts, and that region's own survey coverage value.",
     "widths": [1, 1.2],
     "cols": [
        {"head": "What changes, and what does not", "size": 12, "items": [
            "The arithmetic is identical. A region's count for the chain indicator is divided by that region's survey coverage, and the implied population is carried down the same cascade.",
            "The chain is not re-selected region by region. The one chosen nationally is used everywhere.",
            "The demographic rates are national as well, applied unchanged to every region. This is a known limitation, noted on the next slide."]},
        {"head": "What the approach requires, and what follows", "size": 12, "items": [
            "A region can only be estimated where a survey publishes a coverage value for that region. At present only DHS regional estimates are integrated, so subnational estimation stops at DHS regions; MICS subnational values are not yet used.",
            "Where a region has no survey value for the anchor indicator, the national value is not borrowed as a substitute. That region gets no coverage estimate, and the output shows the gap.",
            "Survey regions are matched to DHIS2 areas by exact name. A region whose name does not match is dropped, and where nothing matches, only national estimates are produced.",
            "A finer second level, such as districts, is possible only where the survey itself publishes estimates at two levels, as the DHS does in Nigeria and Ethiopia.",
            "Estimates above 100% are shown as they are, not capped at 100%. A value above 100% warns that the denominator is wrong for that unit, and for such units the trend is read rather than the level."]},
     ]},

    # Slide 13
    {"type": "table", "title": "Comparison of approaches",
     "intro": "With the FASTR method now described, the two approaches can be set side by side. They differ at almost every step, and the differences follow from the questions each set out to answer.",
     "headers": ["Dimension", "IHME analysis", "FASTR"], "colw": [1.8, 5.15, 5.14], "bold_first": True, "fs": 10,
     "rows": [
        ["What it is", "An evaluation of the administrative data countries have already reported, tested against survey estimates", "A pipeline that constructs estimates for country programs to use"],
        ["Data examined", "eJRF submissions, annual, at first and second administrative level", "DHIS2 extracts, monthly, at facility level"],
        ["Reporting stage", "An official annual submission, compiled and reviewed nationally before it is reported", "Facility records as held in the national system, extracted directly"],
        ["What is examined", "The coverage value, the doses and the target population, all as countries reported them and all used unchanged, since the reported figures are what is under examination", "Facility records assessed for quality and adjusted before use, then divided by a denominator the pipeline builds"],
        ["Denominator", "The target population countries reported, with IHME and WorldPop population surfaces substituted to test how much the choice matters", "No reported population is used; the denominator is derived from service volumes calibrated to survey coverage"],
        ["Survey estimates", "Built from survey microdata for matched birth cohorts, and used as the comparator", "Taken from published survey values, and used as the calibration anchor"],
        ["Indicators", "BCG, DTP1, DTP3 and MCV1", "Reproductive, maternal, newborn, child and adolescent health indicators"],
        ["Question addressed", "Do administrative and survey data agree, and where do they diverge?", "Given survey anchors, what do HMIS trends imply since?"],
     ],
     "note": "These are not competing methods for the same quantity. IHME examines the terms countries produced; FASTR constructs both terms to produce an estimate, so its agreement with the anchoring survey is built in and cannot serve as validation.",
     "notes": "The eJRF and DHIS2 are not the same series, as noted both in Emily's message and in Latera's Ethiopia work. Subnational eJRF data also require geomatching before they can be compared across countries, which the IHME deck credits to WHO and Gavi, and the units reported may not correspond exactly to the DHIS2 hierarchy."},

    # Slide 14
    {"type": "two_col", "title": "Limitations of the FASTR coverage method",
     "widths": [1, 1],
     "cols": [
        {"head": "Dependence on surveys", "size": 11.5, "items": [
            "Coverage cannot be estimated without a survey, and regional coverage needs regional survey estimates. The older the survey, the weaker the estimate, and any bias in it passes straight through.",
            "If our estimates match the survey they were calibrated to, that proves nothing: the agreement was built in. The real test is to anchor to one survey, project forward, then compare against the next survey once it lands. We have not done that yet.",
            "The anchor indicator gets no independent estimate. Its series is produced, but it only returns the survey value, so just the other indicators combine HMIS and survey information.",
            "If the anchor indicator has a data quality problem, every denominator moves with it. All the risk sits in one reporting series."]},
        {"head": "Construction of the denominator", "size": 11.5, "items": [
            "One chain is chosen for the whole country. A chain that fits the country can fit an individual region badly, and that is the likeliest reason a region shows implausible coverage.",
            "The denominator is recomputed every year from that year's anchor count, but the survey coverage it is divided by stays frozen until the next survey, with no expiry rule. If the anchor's coverage really improved in between, the implied population comes out too large. We are exploring options for adding a population growth factor.",
            "The demographic rates are set for each country, but one value is then used for every year and every region. This bears most on subnational estimates, and a study is being scoped to test whether allowing rates to vary by region improves them.",
            "Each rate is itself an estimate, and the errors stack up along the chain. The further an indicator sits from the anchor, the more of that error it carries."]},
     ]},

    # Slide 15
    {"type": "numbered", "title": "Open questions and work under way", "steps": [
        ("Does how we implement this method change the answer?",
         "Our own choices are all open to question: which adjustment scenario feeds the numerator, how the chain is selected, the fixed national cascade rates, the survey value held constant between surveys. If we varied them, would coverage estimates move materially, and would the same units still be identified as lowest?"),
        ("Would a different approach do better?",
         "Comparing our survey-anchored denominator against other sources of population, including census-based targets held in DHIS2, gridded population estimates, and figures a country supplies itself. Does any of them give more valid coverage, in level, in trend, or in both?"),
        ("How would we know?",
         "Both questions need a yardstick. The direct test is to anchor a projection to one survey and compare it against the survey that follows, which is the only comparison our construction does not guarantee."),
     ],
     "footer": "Whether the choice of construction changes which units would be prioritized is the same concern as the bottom-20% analysis in the IHME deck.",
     "notes": "Fuller detail on each thread is in the annex: the methods study in design, the second chain-selection test, extending the benchmarks, out-of-sample validation, and integrating MICS subnational estimates."},

    # Slide 16
    {"type": "section", "title": "Reflections and possible collaboration",
     "kicker": "Response to the questions raised in the thread"},

    # Slide 17
    {"type": "threecol", "title": "Reactions, method limits, and possible joint work",
     "cols": [
        {"head": "Findings consistent with our experience", "items": [
            "A survey-anchored denominator does not make the above-100% problem go away. One chain is chosen nationally, so it can fit a particular region badly, and our own estimates then exceed 100% there. We read that as a wrong denominator for the region, not as coverage.",
            "Aggregation conceals as much as it smooths, as your own discussion notes. A district-year total built from facilities that reported half the year still looks like a normal number, which is why we assess quality at facility-month level.",
            "Trends are more trustworthy than levels, and your work points the same way: numerator-only trends beat coverage trends. But only a little, since agreement with surveys stayed near 50%.",
            "Routine data does catch big shocks. FASTR's disruption analysis picks out strikes, stockouts and the COVID period from the service volumes. Even where the coverage level is unusable, the timing of a disruption shows up."]},
        {"head": "What the adjustments fix, and what they do not", "items": [
            "Cleaning the numerator fills reporting gaps and removes impossible spikes. It cannot help where a facility reports faithfully every month but writes down the wrong number, or counts the wrong service.",
            "Deriving the denominator from service volumes sidesteps the census projection but does not solve the population problem. It trades one dependency for another: an estimate is only as current as the last survey, and exists only where that survey reported.",
            "Whether any of this actually helps is untested. We do not know whether FASTR estimates rank districts, or read direction of change, any better than raw administrative data. The IHME framework is built to answer exactly that.",
            "A better population estimate alone would not fix it. IHME substituted its own and WorldPop surfaces and correlation with surveys barely improved, which points at the numerator as well."]},
        {"head": "Where collaboration could go", "items": [
            "Ten countries in your survey comparison are FASTR countries with subnational estimates: Cameroon, DRC, Ethiopia, Ghana, Guinea, Haiti, Liberia, Madagascar, Nigeria and Sierra Leone. Your evaluation framework could be applied to FASTR-adjusted estimates in any of them.",
            "Latera's Ethiopia comparison of eJRF and DHIS2 data could be extended to further countries, which would establish how far findings from one source carry over to the other. This needs aggregates rather than facility records, so it is the least constrained of these options.",
            "DHIS2 data are held under country agreements, so joint analysis would require country consent and may need to operate on summary measures rather than on facility-level records."]},
     ]},

    # Slide 18
    {"type": "discussion", "title": "Discussion", "questions": [
        "Our outlier and consistency rules differ from the checks in your quality assessment. Where would the two disagree about whether a value is implausible?",
        "Your conclusions concern administrative coverage as countries report it through the eJRF. FASTR builds its series differently: the most recent survey sets the level, and HMIS only moves that level up or down from year to year. Would your findings carry over to a series built that way?",
        "When a new survey arrives, we can compare it against what our projection said beforehand. How many countries, indicators and levels would such a comparison need to convince you?",
        "Would a joint case study in one or two shared countries be of interest, given that DHIS2 data cannot leave a country without its consent?",
     ]},

    # Slide 19
    {"type": "section", "title": "Backup", "kicker": "Reference material"},

    # Slide 20
    {"type": "blocks_cards", "title": "Chain selection in detail", "card_h": 1.15,
     "blocks": [
        {"t": "Cascades are computed from ANC1, institutional delivery, skilled birth attendance and Penta1, all of which are available subnationally, and from BCG at national level. Denominators based on United Nations World Population Prospects estimates are computed nationally to serve as a comparator.", "size": 13.5, "space": 11},
        {"t": "One chain is used for the country as a whole. The pipeline proposes a default: each candidate's implied population is compared with the matching UN WPP estimate across years, and the chain whose median ratio is closest to 1.00 is put forward. The analyst can accept it or impose another on programmatic grounds.", "size": 13.5, "space": 11},
        {"t": "Whichever chain is chosen, it supplies the denominator for every indicator and every geographic level, and denominators from different chains are never mixed. That constraint keeps the implied populations mutually consistent, so that estimated pregnancies exceed live births and live births exceed surviving infants.", "size": 13.5, "space": 0},
     ],
     "cards": [
        ("The anchor indicator",
         "Suppose ANC1 is the anchor: pregnancies are the ANC1 count divided by ANC1 survey coverage. Dividing the ANC1 count back by that denominator just returns the survey value, since the count cancels. The anchor's own series is still produced, but it only returns the survey value, whereas every other indicator yields an independent estimate.", GOLD),
     ],
     "notes": "Where UN WPP data are missing or no years overlap, no chain is proposed automatically and the affected outputs are reported as unavailable. Agreement between chain-based coverage and survey values is computed as a diagnostic but does not enter the proposal."},

    {"type": "numbered", "title": "Worked example: one region's estimate", "steps": [
        ("Take the region's own count",
         "The region recorded 800 first antenatal visits in the year."),
        ("Divide by the region's own survey coverage",
         "The survey put ANC1 coverage in that region at 70%, which implies about 1,140 pregnancies."),
        ("Apply the national demographic rates",
         "The same country-level cascade rates are applied to that figure to obtain the region's other target populations."),
        ("Divide, then extend",
         "Each indicator's regional count is divided by its matching population, and the series is carried past the last survey using the region's own trend."),
     ],
     "footer": "The rates applied here are national. They are not varied by region, which is the limitation stated on the limitations slide."},

    {"type": "numbered", "title": "Open questions in detail", "steps": [
        ("Separating the contribution of the numerator from that of the denominator",
         "A methods study is at the very start of its design. The intent is to build the same coverage under different numerators and denominators, then compare each against survey benchmarks to see which term drives the divergence. Constructions, countries and indicators are not yet chosen. IHME already tried alternative population surfaces for modest gain, so this should extend that rather than repeat it."),
        ("Chain selection",
         "A chain is currently chosen on one test alone: how close its implied populations come to UN WPP. A second test is under evaluation, asking how well the resulting coverage matches survey estimates, since a chain can imply plausible populations and still produce poor coverage. The anchor is excluded from that test, since its coverage equals the survey value regardless."),
        ("Extending the benchmarks",
         "Population targets held in DHIS2 are census-based and available at district level, and together with census counts they would allow subnational fit to enter chain selection. UN WPP cannot support that, since it has no subnational counterpart."),
        ("Out-of-sample validation",
         "The projection assumes that when HMIS coverage moves, real coverage moved too. IHME found administrative and survey data agreeing on the direction of subnational change only about half the time, and dose trends alone did little better. Whether FASTR's anchored trends do better is untested. As new surveys arrive, a projection can be compared against the survey that follows it, which is the direct test."),
        ("Integrating MICS subnational estimates",
         "This would extend subnational anchoring beyond DHS regions, which at present set the limit of where subnational coverage can be estimated."),
     ],
     "footer": "One question the methods study could ask is whether the choice of construction changes which units would be prioritized, which is the same concern as the bottom-20% analysis in the IHME deck, approached by varying the construction rather than by comparing data sources."},

    {"type": "pillars3", "title": "Data sources and their roles", "pillars": [
        {"kicker": "Level", "head": "Household surveys (DHS and MICS)",
         "desc": "These anchor the denominators through back-calculation at the chain entry point, and subnational anchoring relies on the regional estimates a survey publishes. Where no survey covers a year, the reference series used for comparison and for projection draws on other published estimates, including those of WHO and UNICEF for immunization."},
        {"kicker": "Trend", "head": "Routine health information (DHIS2)",
         "desc": "This provides the adjusted numerators, the chain-indicator volumes that carry the denominator across years and regions, and the year-on-year change that extends coverage beyond the most recent survey."},
        {"kicker": "Comparator", "head": "UN World Population Prospects",
         "desc": "This provides the benchmark against which candidate chains are judged during selection, and it never supplies the denominator used to calculate coverage."},
     ]},

    # Slide 21
    {"type": "table", "title": "Demographic cascade rates",
     "intro": "The values below are examples only. Each rate is set for the country being analyzed, and once set it is applied to every year and every administrative unit in that analysis.",
     "headers": ["Parameter", "Example value", "Step in the cascade"], "colw": [3.4, 1.8, 6.89], "bold_first": True, "fs": 11.5,
     "rows": [
        ["Pregnancy loss", "3%", "Pregnancies to deliveries"],
        ["Twin births", "1.5%", "Deliveries to births"],
        ["Stillbirth", "2%", "Births to live births"],
        ["Neonatal mortality", "3.9%", "Live births to DPT-eligible infants"],
        ["Post-neonatal mortality", "2.8%", "DPT-eligible infants to measles-eligible children"],
     ]},

    # Slide 22
    {"type": "contrast", "title": "External benchmarking against Countdown to 2030",
     "cols": [
        {"head": "The same raw data through two pipelines", "accent": DEEP_GREEN, "items": [
            "The same raw DHIS2 data were run through the FASTR and Countdown to 2030 pipelines for seven countries in 2022, for institutional live births and cesarean deliveries.",
            "Countdown adjusts at district level, using a threshold of five median absolute deviations and a k-factor that scales reported totals to allow for incomplete reporting.",
            "FASTR adjusts at facility level, using a threshold of ten deviations and imputing directly from each facility's own reporting history."]},
        {"head": "What the comparison found", "accent": GREEN, "items": [
            "The median relative difference in adjusted national totals was **0.9%** for institutional live births and **3.3%** for cesarean deliveries.",
            "At first administrative level, median differences fell below 2% in five of the six countries with subnational comparisons.",
            "Malawi was the exception, at 6.2% for live births and 29.7% for cesarean deliveries, and it has been identified for joint investigation with the Countdown team."]},
     ],
     "bottom": [
        ("What follows from it",
         "Two independently developed adjustment methods produce comparable national numerators in six of seven countries, despite differing in the unit at which they adjust, in their detection threshold and in their treatment of incomplete reporting.", DEEP_GREEN),
     ],
     "notes": "This is held in reserve rather than presented, and it answers the question of whether the adjustment has been checked against anything external. A structured comparison, including a decomposition of where the two pipelines diverge, was shared with the Countdown team together with questions still open for them."},

    # Slide 23
    {"type": "table", "title": "Adjustment methods: FASTR and Countdown to 2030",
     "intro": "Drawn from the comparison working document of March 2026, which was shared with the Countdown team.",
     "headers": ["Dimension", "FASTR", "Countdown to 2030"], "colw": [2.4, 5.0, 4.69], "bold_first": True, "fs": 11,
     "rows": [
        ["Unit of adjustment", "Facility-month", "District-month"],
        ["Outlier rule", "Ten median absolute deviations from the facility median, or a single month exceeding 80% of trailing twelve-month volume, applied to counts above one hundred", "Five median absolute deviations"],
        ["Incomplete reporting", "Imputed directly from each facility's reporting history, through a fixed hierarchy", "Reported totals scaled by a k-factor, with k set to 0.25"],
        ["Sensitivity analysis", "Four datasets carried through the pipeline", "Adjustment factors applied to reported aggregates"],
     ]},
]


def main():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    for c in CONTENT:
        BUILDERS[c["type"]](prs, c)
    total = len(CONTENT)
    no_footer = {"cover", "section", "discussion"}
    for n, (c, slide) in enumerate(zip(CONTENT, prs.slides), 1):
        _fade(slide)
        if c["type"] not in no_footer:
            _footer(slide, n, total)
        if c.get("notes"):
            slide.notes_slide.notes_text_frame.text = c["notes"]
    out = REPO / "decks" / "en" / "fastr_ihme_methods.pptx"
    prs.save(str(out))
    print(f"Wrote {len(CONTENT)} slides -> {out}")


if __name__ == "__main__":
    main()
