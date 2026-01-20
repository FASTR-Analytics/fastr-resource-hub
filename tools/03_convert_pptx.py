#!/usr/bin/env python3
"""
═══════════════════════════════════════════════════════════════════════
              FASTR POWERPOINT CONVERTER
═══════════════════════════════════════════════════════════════════════

Converts Marp markdown decks to branded PowerPoint presentations.

USAGE:
    python3 tools/03_convert_pptx.py                           # Interactive
    python3 tools/03_convert_pptx.py outputs/workshop_deck.md  # Direct

FEATURES:
  ✓ FASTR brand colors and styling
  ✓ Proper slide layouts (title, content, agenda, two-column)
  ✓ Tables with styled headers
  ✓ Images embedded correctly
  ✓ Break slides with centered text

═══════════════════════════════════════════════════════════════════════
"""

import argparse
import os
import re
import sys
from pathlib import Path

# ═══════════════════════════════════════════════════════════════════════════════
# AUTO-DETECT AND USE VENV
# ═══════════════════════════════════════════════════════════════════════════════

def ensure_venv():
    """Re-execute with venv Python if not already in venv."""
    if sys.prefix != sys.base_prefix:
        return
    script_dir = Path(__file__).resolve().parent
    project_root = script_dir.parent
    for venv_name in ['.venv', 'venv']:
        venv_python = project_root / venv_name / 'bin' / 'python3'
        if venv_python.exists():
            os.execv(str(venv_python), [str(venv_python)] + sys.argv)

ensure_venv()

# Now import python-pptx (after venv activation)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed.")
    print("Run: pip install python-pptx Pillow")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    Image = None

# ═══════════════════════════════════════════════════════════════════════════════
# FASTR BRAND CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════

class Colors:
    """FASTR brand colors from fastr-theme.css"""
    DEEP_GREEN = RGBColor(0x09, 0x54, 0x4F)   # #09544F - H1
    DARK_GREEN = RGBColor(0x0C, 0x71, 0x6B)   # #0C716B
    GREEN = RGBColor(0x1F, 0x9A, 0x9C)        # #1F9A9C
    LIME = RGBColor(0xD0, 0xCB, 0x17)         # #D0CB17 - H1 underline
    NAVY = RGBColor(0x21, 0x56, 0x8C)         # #21568C - H2
    BLUE = RGBColor(0x1A, 0x90, 0xC0)         # #1A90C0 - H2 underline
    LIGHT_BLUE = RGBColor(0xCA, 0xE6, 0xE9)   # #CAE6E9 - table headers
    LIGHT_GREEN = RGBColor(0xE8, 0xF4, 0xF3)  # #E8F4F3 - session headers

    GOLD = RGBColor(0xD8, 0xA8, 0x22)         # #D8A822
    PURPLE = RGBColor(0x7A, 0x1F, 0x6E)       # #7A1F6E
    ORCHID = RGBColor(0xBD, 0x50, 0x91)       # #BD5091
    CORAL = RGBColor(0xFF, 0x64, 0x62)        # #FF6462

    TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)    # #2c3e50
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)    # #333333 - body text
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x00, 0x00, 0x00)


class Fonts:
    """Font settings"""
    FAMILY = 'Poppins'  # FASTR brand typeface
    H1_SIZE = Pt(32)   # Reduced from Pt(40) to prevent underline overlap
    H2_SIZE = Pt(26)   # Reduced from Pt(32) to prevent underline overlap
    H3_SIZE = Pt(22)
    BODY_SIZE = Pt(18)  # Increased from 16 for better readability
    TABLE_SIZE = Pt(14)
    SMALL_SIZE = Pt(12)


class Layout:
    """Slide dimensions and margins (16:9)"""
    WIDTH = Inches(13.333)
    HEIGHT = Inches(7.5)
    MARGIN_LEFT = Inches(0.75)
    MARGIN_RIGHT = Inches(0.75)
    MARGIN_TOP = Inches(0.6)
    MARGIN_BOTTOM = Inches(0.5)
    CONTENT_WIDTH = Inches(11.833)  # 13.333 - 0.75 - 0.75
    # Centered content position: (slide_width - content_width) / 2
    CONTENT_LEFT = Inches(0.75)


# ═══════════════════════════════════════════════════════════════════════════════
# MARKDOWN PARSER
# ═══════════════════════════════════════════════════════════════════════════════

def parse_markdown(content):
    """
    Parse Marp markdown into slide data structures.

    Returns list of dicts with:
    - raw: original markdown
    - headers: [(level, text), ...]
    - bullets: [text, ...]
    - table: [[row], ...] or None
    - images: [{'alt': str, 'path': str}, ...]
    - columns: {'left': str, 'right': str} or None
    - css_class: from <!-- _class: xxx -->
    """
    # Strip YAML frontmatter
    if content.startswith('---'):
        match = re.match(r'^---\n.*?\n---\n?', content, re.DOTALL)
        if match:
            content = content[match.end():]

    # Strip <style> blocks (including scoped)
    content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL)

    # Split into slides
    raw_slides = re.split(r'\n---\s*\n', content)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        slide = {
            'raw': raw,
            'headers': [],
            'bullets': [],
            'paragraphs': [],  # Plain text lines (not headers, bullets, or tables)
            'content': [],     # Ordered list of ('paragraph', text) or ('bullet', text)
            'table': None,
            'images': [],
            'columns': None,
            'css_class': None,
            'html_content': None,
        }

        # Extract CSS class directive
        class_match = re.search(r'<!--\s*_class:\s*(\w+)\s*-->', raw)
        if class_match:
            slide['css_class'] = class_match.group(1)
            raw = re.sub(r'<!--\s*_class:\s*\w+\s*-->', '', raw)

        # Extract headers (strip HTML tags like <img>)
        for match in re.finditer(r'^(#{1,6})\s+(.+)$', raw, re.MULTILINE):
            level = len(match.group(1))
            text = match.group(2).strip()
            # Strip HTML tags from headers
            text = re.sub(r'<[^>]+>', '', text).strip()
            slide['headers'].append((level, text))

        # Extract images ![alt](path)
        for match in re.finditer(r'!\[([^\]]*)\]\(([^)]+)\)', raw):
            slide['images'].append({
                'alt': match.group(1),
                'path': match.group(2).split()[0]  # Remove any title
            })

        # Extract HTML images <img src="path">
        for match in re.finditer(r'<img\s+[^>]*src=["\']([^"\']+)["\']', raw):
            slide['images'].append({
                'alt': '',
                'path': match.group(1)
            })

        # Extract table (markdown or HTML)
        table_lines = []
        in_table = False

        # First try HTML table
        html_table_match = re.search(r'<table[^>]*>(.*?)</table>', raw, re.DOTALL | re.IGNORECASE)
        if html_table_match:
            table_html = html_table_match.group(1)
            # Parse rows
            for row_match in re.finditer(r'<tr[^>]*>(.*?)</tr>', table_html, re.DOTALL | re.IGNORECASE):
                row_html = row_match.group(1)
                cells = []
                # Parse cells (th or td)
                for cell_match in re.finditer(r'<t[hd][^>]*>(.*?)</t[hd]>', row_html, re.DOTALL | re.IGNORECASE):
                    cell_text = cell_match.group(1)
                    # Check for colspan
                    colspan_match = re.search(r'colspan=["\']?(\d+)', row_match.group(0), re.IGNORECASE)
                    # Strip HTML tags from cell content but preserve text
                    cell_text = re.sub(r'<strong>|</strong>', '**', cell_text)
                    cell_text = re.sub(r'<em>|</em>', '*', cell_text)
                    cell_text = re.sub(r'<[^>]+>', '', cell_text).strip()
                    cells.append(cell_text)
                    # If colspan, add empty cells
                    if colspan_match:
                        for _ in range(int(colspan_match.group(1)) - 1):
                            cells.append('')
                if cells:
                    table_lines.append(cells)
        else:
            # Try markdown table
            for line in raw.split('\n'):
                if '|' in line and line.strip().startswith('|'):
                    in_table = True
                    # Skip separator line
                    if re.match(r'^\|[\s\-:|]+\|$', line.strip()):
                        continue
                    cells = [c.strip() for c in line.strip().strip('|').split('|')]
                    table_lines.append(cells)
                elif in_table:
                    break

        if table_lines:
            slide['table'] = table_lines

        # Extract columns div
        cols_match = re.search(
            r'<div\s+class="columns[^"]*">\s*<div>(.*?)</div>\s*<div>(.*?)</div>\s*</div>',
            raw, re.DOTALL
        )
        if cols_match:
            slide['columns'] = {
                'left': cols_match.group(1).strip(),
                'right': cols_match.group(2).strip()
            }

        # Check for other HTML content (like styled img tags)
        if '<img' in raw:
            slide['html_content'] = raw

        # Extract content IN ORDER - headers (H3+), bullets, numbered lists, and paragraphs
        for line in raw.split('\n'):
            line_stripped = line.strip()
            if not line_stripped:
                continue
            # Include H3+ headers in content flow (H1/H2 rendered separately at top)
            header_match = re.match(r'^(#{3,6})\s+(.+)$', line_stripped)
            if header_match:
                level = len(header_match.group(1))
                text = header_match.group(2).strip()
                slide['content'].append(('header', (level, text)))
                continue
            # Skip H1/H2 headers (rendered at top)
            if line_stripped.startswith('#'):
                continue
            # Skip table lines
            if line_stripped.startswith('|'):
                continue
            # Skip images (markdown and HTML)
            if line_stripped.startswith('!['):
                continue
            if line_stripped.startswith('<img'):
                continue
            # Skip div wrappers (but HTML text tags like <em> handled later)
            if line_stripped.startswith('<div') or line_stripped.startswith('</div'):
                continue
            # Skip comments
            if line_stripped.startswith('<!--'):
                continue

            # Check for bullet
            bullet_match = re.match(r'^[-*]\s+(.+)$', line_stripped)
            if bullet_match:
                text = bullet_match.group(1).strip()
                slide['bullets'].append(text)
                slide['content'].append(('bullet', text))
                continue

            # Check for numbered list
            num_match = re.match(r'^\d+\.\s+(.+)$', line_stripped)
            if num_match:
                text = num_match.group(1).strip()
                slide['bullets'].append(text)
                slide['content'].append(('bullet', text))
                continue

            # It's a paragraph
            slide['paragraphs'].append(line_stripped)
            slide['content'].append(('paragraph', line_stripped))

        slides.append(slide)

    return slides


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE TYPE DETECTOR
# ═══════════════════════════════════════════════════════════════════════════════

def detect_slide_type(slide, index):
    """
    Detect slide type from content.

    Returns: 'title', 'agenda', 'break', 'section', 'two_column', 'table', 'image', 'content'
    """
    headers = slide['headers']
    h1_text = headers[0][1] if headers and headers[0][0] == 1 else ''

    # Title slide: first slide with logo
    if index == 0:
        for img in slide['images']:
            if 'logo' in img['path'].lower() or 'fastr' in img['path'].lower():
                return 'title'

    # Agenda slide
    if slide['css_class'] == 'agenda' or (slide['table'] and 'agenda' in h1_text.lower()):
        return 'agenda'

    # Break slide: emoji or "break" in H1
    break_emojis = ['☕', '🍽', '🌙', '🎉', '👋', '⏰']
    if any(emoji in h1_text for emoji in break_emojis):
        return 'break'
    if re.search(r'\b(break|lunch|tea)\b', h1_text, re.IGNORECASE):
        return 'break'

    # Section header slide: "Session X:" ONLY when minimal content (just header)
    # Don't use section type if slide has bullets, paragraphs, or images
    has_content = len(slide['bullets']) > 0 or len(slide['paragraphs']) > 0 or slide['images']
    if headers and headers[0][0] == 1 and not has_content:
        if re.match(r'^Session\s+\d+', h1_text, re.IGNORECASE):
            return 'section'

    # Two-column slide
    if slide['columns']:
        return 'two_column'

    # Table slide (non-agenda)
    if slide['table']:
        return 'table'

    # Image-heavy slide
    if slide['images'] and len(slide['bullets']) < 2:
        return 'image'

    # Default: content slide
    return 'content'


# ═══════════════════════════════════════════════════════════════════════════════
# IMAGE HANDLING
# ═══════════════════════════════════════════════════════════════════════════════

def resolve_image_path(img_path, base_dir, md_dir):
    """
    Resolve relative image path to absolute path.
    Tries multiple locations.
    """
    if img_path.startswith(('http://', 'https://')):
        return None  # Skip URLs

    # Clean up the path
    img_path = img_path.replace('%20', ' ')

    paths_to_try = [
        os.path.join(md_dir, img_path),
        os.path.join(base_dir, img_path.lstrip('../')),
        os.path.join(base_dir, 'resources', 'logos', os.path.basename(img_path)),
        os.path.join(base_dir, 'resources', 'diagrams', os.path.basename(img_path)),
        os.path.join(base_dir, 'resources', 'default_outputs', os.path.basename(img_path)),
        os.path.join(base_dir, 'assets', os.path.basename(img_path)),
    ]

    # Handle ../ paths from outputs/ folder (resources, workshops, etc.)
    if img_path.startswith('../'):
        clean_path = img_path.replace('../', '')
        paths_to_try.insert(0, os.path.join(base_dir, clean_path))

    for path in paths_to_try:
        if os.path.exists(path):
            return os.path.abspath(path)

    return None


def convert_svg_to_png(svg_path, temp_dir):
    """Convert SVG to PNG for PowerPoint compatibility."""
    try:
        import cairosvg
        png_path = os.path.join(temp_dir, Path(svg_path).stem + '.png')
        cairosvg.svg2png(url=svg_path, write_to=png_path, scale=2.0)
        return png_path
    except ImportError:
        return None
    except Exception:
        return None


def add_image_to_slide(slide, img_path, left, top, width=None, height=None):
    """Add image to slide, handling SVG conversion if needed."""
    if not img_path or not os.path.exists(img_path):
        print(f"   Warning: Image not found: {img_path}")
        return None

    # Convert SVG if needed
    if img_path.lower().endswith('.svg'):
        import tempfile
        # Create a persistent temp file (not auto-deleted)
        temp_dir = tempfile.mkdtemp()
        png_path = convert_svg_to_png(img_path, temp_dir)
        if png_path and os.path.exists(png_path):
            img_path = png_path
        else:
            print(f"   Warning: SVG conversion failed for {os.path.basename(img_path)} - install cairosvg")
            return None

    try:
        if width and height:
            return slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            return slide.shapes.add_picture(img_path, left, top)
    except Exception as e:
        print(f"   Warning: Could not add image {os.path.basename(img_path)}: {e}")
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# TEXT STYLING HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def style_paragraph(para, font_size, color, bold=False, italic=False):
    """Apply styling to a paragraph."""
    para.font.size = font_size
    para.font.color.rgb = color
    para.font.bold = bold
    para.font.italic = italic
    para.font.name = Fonts.FAMILY


def strip_html_tags(text):
    """Remove HTML tags from text, converting em/i to markdown italic first."""
    # Convert <em>...</em> and <i>...</i> to markdown italic *...*
    text = re.sub(r'<em>([^<]+)</em>', r'*\1*', text, flags=re.IGNORECASE)
    text = re.sub(r'<i>([^<]+)</i>', r'*\1*', text, flags=re.IGNORECASE)
    # Convert <strong>...</strong> and <b>...</b> to markdown bold **...**
    text = re.sub(r'<strong>([^<]+)</strong>', r'**\1**', text, flags=re.IGNORECASE)
    text = re.sub(r'<b>([^<]+)</b>', r'**\1**', text, flags=re.IGNORECASE)
    # Convert <small>...</small> to regular text (italicized)
    text = re.sub(r'<small>([^<]+)</small>', r'*\1*', text, flags=re.IGNORECASE)
    # Remove <img ...> tags
    text = re.sub(r'<img\s+[^>]*/?>', '', text)
    # Remove other common HTML tags but keep content
    text = re.sub(r'</?(?:div|span|p|br|a|small|em|i|strong|b)[^>]*>', '', text)
    return text.strip()


def strip_markdown_images(text):
    """Remove markdown image syntax from text."""
    return re.sub(r'!\[[^\]]*\]\([^)]+\)', '', text).strip()


def add_formatted_text(paragraph, text, font_size, color):
    """
    Parse markdown formatting and add styled runs to paragraph.
    Handles **bold**, *italic*, and `code` formatting.
    Bold text is colored Navy, italic text is colored Green (matching CSS).
    """
    # Clear any existing text
    paragraph.clear()

    # Strip HTML tags first (icons show as <img...> otherwise)
    text = strip_html_tags(text)

    # Pattern to match **bold**, *italic*, `code`, or plain text
    # Use non-greedy matching
    pattern = r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`|[^*`]+)'

    for match in re.finditer(pattern, text):
        segment = match.group(1)
        if not segment:
            continue

        run = paragraph.add_run()
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = Fonts.FAMILY

        if segment.startswith('**') and segment.endswith('**'):
            # Bold text - Navy color to match CSS
            run.text = segment[2:-2]
            run.font.bold = True
            run.font.color.rgb = Colors.NAVY
        elif segment.startswith('*') and segment.endswith('*') and not segment.startswith('**'):
            # Italic text - Green color to match CSS
            run.text = segment[1:-1]
            run.font.italic = True
            run.font.color.rgb = Colors.GREEN
        elif segment.startswith('`') and segment.endswith('`'):
            # Code text - Navy color
            run.text = segment[1:-1]
            run.font.name = 'Courier New'
            run.font.color.rgb = Colors.NAVY
        else:
            # Plain text
            run.text = segment


def add_formatted_bullet_text(paragraph, text, font_size, color):
    """
    Add a bullet point with lime-colored bullet and formatted text.
    The bullet character is rendered in lime, the text uses standard formatting.
    """
    # Clear any existing text
    paragraph.clear()

    # Add lime-colored bullet character
    bullet_run = paragraph.add_run()
    bullet_run.text = "• "
    bullet_run.font.size = font_size
    bullet_run.font.color.rgb = Colors.LIME
    bullet_run.font.name = Fonts.FAMILY

    # Strip HTML tags first
    text = strip_html_tags(text)

    # Pattern to match **bold**, *italic*, `code`, or plain text
    pattern = r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`|[^*`]+)'

    for match in re.finditer(pattern, text):
        segment = match.group(1)
        if not segment:
            continue

        run = paragraph.add_run()
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = Fonts.FAMILY

        if segment.startswith('**') and segment.endswith('**'):
            # Bold text - Navy color
            run.text = segment[2:-2]
            run.font.bold = True
            run.font.color.rgb = Colors.NAVY
        elif segment.startswith('*') and segment.endswith('*') and not segment.startswith('**'):
            # Italic text - Green color
            run.text = segment[1:-1]
            run.font.italic = True
            run.font.color.rgb = Colors.GREEN
        elif segment.startswith('`') and segment.endswith('`'):
            # Code text - Navy color
            run.text = segment[1:-1]
            run.font.name = 'Courier New'
            run.font.color.rgb = Colors.NAVY
        else:
            # Plain text
            run.text = segment


def add_text_box(slide, left, top, width, height, text, font_size, color,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a styled text box to slide."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Set vertical alignment
    try:
        tf.anchor = anchor
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    style_paragraph(p, font_size, color, bold)

    return shape


def add_h1_with_underline(slide, text, top=None):
    """Add H1 header with lime underline."""
    if top is None:
        top = Layout.MARGIN_TOP

    # Add title text - centered
    title_shape = add_text_box(
        slide,
        Layout.CONTENT_LEFT, top,
        Layout.CONTENT_WIDTH, Inches(0.6),
        text, Fonts.H1_SIZE, Colors.DEEP_GREEN, bold=True
    )

    # Add underline - positioned well below the text
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Layout.CONTENT_LEFT, top + Inches(0.65),
        Inches(4), Inches(0.04)  # Thicker to match PDF styling
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = Colors.LIME
    underline.line.fill.background()

    return title_shape


def add_h2_with_underline(slide, text, top):
    """Add H2 header with blue underline."""
    title_shape = add_text_box(
        slide,
        Layout.CONTENT_LEFT, top,
        Layout.CONTENT_WIDTH, Inches(0.5),
        text, Fonts.H2_SIZE, Colors.NAVY, bold=True
    )

    # Add underline - positioned well below the text
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Layout.CONTENT_LEFT, top + Inches(0.55),
        Inches(3), Inches(0.03)  # Thicker to match PDF styling
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = Colors.BLUE
    underline.line.fill.background()

    return title_shape


def add_bullet_list(slide, bullets, left, top, width, font_size=None):
    """Add a bullet list to slide with proper markdown formatting."""
    if not bullets:
        return None

    if font_size is None:
        font_size = Fonts.BODY_SIZE

    height = Inches(0.35 * len(bullets) + 0.2)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Use formatted bullet text with lime-colored bullet
        add_formatted_bullet_text(p, bullet, font_size, Colors.TEXT_DARK)
        p.level = 0

    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_title_slide(prs, data, base_dir, md_dir):
    """Build title slide with background image and white text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Check for background image (![bg] syntax in raw)
    raw = data['raw']
    bg_match = re.search(r'!\[bg\]\(([^)]+)\)', raw)
    has_bg_image = False

    if bg_match:
        bg_path = bg_match.group(1).split()[0]
        bg_full_path = resolve_image_path(bg_path, base_dir, md_dir)
        if bg_full_path and os.path.exists(bg_full_path):
            # Set background image to fill entire slide
            slide.shapes.add_picture(
                bg_full_path,
                Inches(0), Inches(0),
                width=Layout.WIDTH, height=Layout.HEIGHT
            )
            has_bg_image = True

    # Get title
    title = data['headers'][0][1] if data['headers'] else 'FASTR Workshop'

    # Text colors - white if background image, otherwise branded colors
    title_color = Colors.WHITE if has_bg_image else Colors.DEEP_GREEN
    subtitle_color = RGBColor(0xF0, 0xF0, 0xF0) if has_bg_image else Colors.NAVY
    facilitator_color = RGBColor(0xE0, 0xE0, 0xE0) if has_bg_image else Colors.TEXT_DARK

    # Centered title
    add_text_box(
        slide,
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5),
        title, Fonts.H1_SIZE, title_color,
        bold=True, align=PP_ALIGN.CENTER
    )

    # Underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3), Inches(4.0),
        Inches(7.333), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = Colors.LIME
    underline.line.fill.background()

    # Look for date/location line (bold text)
    date_match = re.search(r'\*\*([^*]+)\*\*\s*\|\s*\*\*([^*]+)\*\*', raw)
    if date_match:
        subtitle = f"{date_match.group(1)} | {date_match.group(2)}"
        add_text_box(
            slide,
            Inches(1), Inches(4.3),
            Inches(11.333), Inches(0.5),
            subtitle, Fonts.H2_SIZE, subtitle_color,
            align=PP_ALIGN.CENTER
        )

    # Look for facilitator (italic text)
    fac_match = re.search(r'\*([^*]+)\*(?!\*)', raw)
    if fac_match:
        add_text_box(
            slide,
            Inches(1), Inches(5.0),
            Inches(11.333), Inches(0.4),
            fac_match.group(1), Fonts.BODY_SIZE, facilitator_color,
            align=PP_ALIGN.CENTER
        )

    # Only add logo if no background image (background already has logos)
    if not has_bg_image:
        for img in data['images']:
            img_path = resolve_image_path(img['path'], base_dir, md_dir)
            if img_path and ('logo' in img_path.lower() or 'fastr' in img_path.lower()):
                add_image_to_slide(
                    slide, img_path,
                    Inches(10.5), Inches(6.2),
                    width=Inches(2.3)
                )
                break

    return slide


def build_agenda_slide(prs, data, base_dir, md_dir):
    """Build agenda slide with compact table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    title = data['headers'][0][1] if data['headers'] else 'Workshop Agenda'
    add_h1_with_underline(slide, title)

    # Table
    if data['table']:
        table_data = data['table']
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 3

        # Use smaller font for agendas with many rows
        font_size = Fonts.SMALL_SIZE if rows > 10 else Fonts.TABLE_SIZE
        row_height = Inches(0.28) if rows > 10 else Inches(0.32)

        # Full width table respecting margins
        table_width = Layout.CONTENT_WIDTH  # 11.833"
        table_left = Layout.CONTENT_LEFT    # 0.75"

        table = slide.shapes.add_table(
            rows, cols,
            table_left, Inches(1.2),
            table_width, row_height * rows
        ).table

        # Set column widths - Time | Session | Speaker (proportional to table width)
        if cols >= 3:
            table.columns[0].width = Inches(2.0)   # Time
            table.columns[1].width = Inches(7.6)   # Session
            table.columns[2].width = Inches(2.2)   # Speaker

        # Fill table with formatted text
        for r_idx, row in enumerate(table_data):
            is_header = r_idx == 0

            # Check if this is a session header row (Session X: or **Header** with empty other cells)
            row_text = ' '.join(str(cell) for cell in row)
            # Match "Session X:" or rows where first cell is bold and other cells are empty
            is_session_row = bool(re.match(r'.*Session\s+\d+:', row_text, re.IGNORECASE))
            # Also check for bold header rows like "**Opening Session**" with empty subsequent cells
            if not is_session_row and cols >= 2:
                first_cell_bold = row[0].startswith('**') and row[0].endswith('**')
                other_cells_empty = all(str(c).strip() == '' for c in row[1:])
                is_session_row = first_cell_bold and other_cells_empty

            # Merge cells for session header rows
            if is_session_row and cols >= 2:
                # Merge all cells in this row
                first_cell = table.cell(r_idx, 0)
                last_cell = table.cell(r_idx, cols - 1)
                first_cell.merge(last_cell)

                # Style the merged cell
                merged_cell = table.cell(r_idx, 0)
                merged_cell.fill.solid()
                merged_cell.fill.fore_color.rgb = Colors.LIGHT_GREEN
                merged_cell.text_frame.word_wrap = False

                para = merged_cell.text_frame.paragraphs[0]
                # Get the session text (find the cell with "Session")
                session_text = next((str(c) for c in row if 'Session' in str(c)), row_text)
                clean_text = re.sub(r'\*+([^*]+)\*+', r'\1', session_text)
                para.text = clean_text.strip()
                para.font.size = font_size
                para.font.name = Fonts.FAMILY
                para.font.color.rgb = Colors.DEEP_GREEN
                para.font.bold = True
            else:
                for c_idx, cell_text in enumerate(row):
                    if c_idx >= cols:
                        continue
                    cell = table.cell(r_idx, c_idx)
                    cell.text_frame.word_wrap = True

                    para = cell.text_frame.paragraphs[0]

                    if is_header:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = Colors.LIGHT_BLUE
                        clean_text = re.sub(r'\*+([^*]+)\*+', r'\1', cell_text)
                        para.text = clean_text
                        para.font.size = font_size
                        para.font.name = Fonts.FAMILY
                        para.font.color.rgb = Colors.NAVY
                        para.font.bold = True
                    else:
                        add_formatted_text(para, cell_text, font_size, Colors.TEXT_DARK)

    return slide


def build_break_slide(prs, data, base_dir, md_dir):
    """Build break slide with large centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Get title (with emoji)
    title = data['headers'][0][1] if data['headers'] else 'Break'

    # Large centered title
    add_text_box(
        slide,
        Inches(1), Inches(2.8),
        Inches(11.333), Inches(1.5),
        title, Pt(56), Colors.DEEP_GREEN,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )

    # Look for time info in raw content
    raw = data['raw']
    time_match = re.search(r'(\d+)\s*minutes?', raw, re.IGNORECASE)
    resume_match = re.search(r'(?:resume|back|return)[^\d]*(\d+[:\d]*\s*(?:AM|PM)?)', raw, re.IGNORECASE)

    subtitle_parts = []
    if time_match:
        subtitle_parts.append(f"{time_match.group(1)} minutes")
    if resume_match:
        subtitle_parts.append(f"Back at {resume_match.group(1)}")

    if subtitle_parts:
        add_text_box(
            slide,
            Inches(1), Inches(4.3),
            Inches(11.333), Inches(0.6),
            " • ".join(subtitle_parts), Fonts.H2_SIZE, Colors.NAVY,
            align=PP_ALIGN.CENTER
        )

    return slide


def build_section_slide(prs, data, base_dir, md_dir):
    """Build section header slide with centered title (like break but for sessions)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Get title
    title = data['headers'][0][1] if data['headers'] else 'Section'

    # Large centered title
    add_text_box(
        slide,
        Inches(1), Inches(2.8),
        Inches(11.333), Inches(1.5),
        title, Pt(44), Colors.DEEP_GREEN,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )

    # Add underline below title
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(4.5),
        Inches(5.333), Inches(0.03)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = Colors.LIME
    underline.line.fill.background()

    return slide


def parse_column_content(content):
    """Parse column content into structured items (paragraphs, bullets, numbered)."""
    items = []
    lines = content.strip().split('\n')

    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Skip image tags only (not text with HTML formatting)
        if line.startswith('!['):
            continue
        if line.startswith('<img'):
            continue
        # Skip div wrappers but not content inside
        if line.startswith('<div') or line.startswith('</div'):
            continue

        # Check for bullet
        bullet_match = re.match(r'^[-*]\s+(.+)$', line)
        if bullet_match:
            items.append(('bullet', bullet_match.group(1)))
            continue

        # Check for numbered list
        num_match = re.match(r'^\d+\.\s+(.+)$', line)
        if num_match:
            items.append(('bullet', num_match.group(1)))
            continue

        # Regular paragraph (may contain HTML tags like <em>, which will be converted)
        items.append(('paragraph', line))

    return items


def add_column_text(slide, items, left, top, width):
    """Add column text content to slide."""
    if not items:
        return None

    # Use full available height
    height = Inches(5.5)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.MIDDLE  # Vertically center text

    first_para = True
    for item_type, text in items:
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        if item_type == 'bullet':
            add_formatted_bullet_text(p, text, Fonts.BODY_SIZE, Colors.TEXT_DARK)
            p.space_after = Pt(8)
        else:
            add_formatted_text(p, text, Fonts.BODY_SIZE, Colors.DARK_GRAY)
            p.space_after = Pt(12)  # More space after paragraphs

    return shape


def build_two_column_slide(prs, data, base_dir, md_dir):
    """Build two-column slide with text and image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    if data['headers']:
        title = data['headers'][0][1]
        if data['headers'][0][0] == 1:
            add_h1_with_underline(slide, title)
        else:
            add_h2_with_underline(slide, title, Layout.MARGIN_TOP)

    # Parse column content
    left_content = data['columns']['left']
    right_content = data['columns']['right']

    # Determine which side has image
    left_has_image = '![' in left_content or '<img' in left_content
    right_has_image = '![' in right_content or '<img' in right_content

    content_top = Inches(1.4)
    col_width = Inches(5.8)

    def extract_image_path(content):
        """Extract image path from markdown or HTML img tag."""
        # Try markdown format first
        md_match = re.search(r'!\[[^\]]*\]\(([^)]+)\)', content)
        if md_match:
            return md_match.group(1).split()[0]
        # Try HTML img tag
        html_match = re.search(r'<img\s+[^>]*src=["\']([^"\']+)["\']', content)
        if html_match:
            return html_match.group(1)
        return None

    # Left column
    if left_has_image:
        # Add image
        img_src = extract_image_path(left_content)
        if img_src:
            img_path = resolve_image_path(img_src, base_dir, md_dir)
            if img_path:
                add_image_to_slide(
                    slide, img_path,
                    Layout.MARGIN_LEFT, content_top,
                    width=col_width
                )
    else:
        # Add text content (paragraphs, bullets, numbered lists)
        items = parse_column_content(left_content)
        if items:
            add_column_text(slide, items, Layout.MARGIN_LEFT, content_top, col_width)

    # Right column
    right_left = Inches(7)
    if right_has_image:
        img_src = extract_image_path(right_content)
        if img_src:
            img_path = resolve_image_path(img_src, base_dir, md_dir)
            if img_path:
                add_image_to_slide(
                    slide, img_path,
                    right_left, content_top,
                    width=col_width
                )
    else:
        # Add text content
        items = parse_column_content(right_content)
        if items:
            add_column_text(slide, items, right_left, content_top, col_width)

    return slide


def build_table_slide(prs, data, base_dir, md_dir):
    """Build slide with table - auto-fit columns based on content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    if data['headers']:
        title = data['headers'][0][1]
        if data['headers'][0][0] == 1:
            add_h1_with_underline(slide, title)
        else:
            add_h2_with_underline(slide, title, Layout.MARGIN_TOP)

    # Check if this is a layout table (contains images) - render differently
    if data['images'] and data['table']:
        table_data = data['table']
        # Check if any table cell contains image markdown
        has_image_cells = any('![' in str(cell) for row in table_data for cell in row)

        if has_image_cells:
            # This is a layout table - render images and text separately
            num_images = len(data['images'])

            # Render images side by side
            if num_images >= 2:
                img_width = Inches(5.5)
                img_top = Inches(1.3)
                for i, img in enumerate(data['images'][:2]):
                    img_path = resolve_image_path(img['path'], base_dir, md_dir)
                    if img_path:
                        left = Layout.CONTENT_LEFT if i == 0 else Inches(7)
                        add_image_to_slide(slide, img_path, left, img_top, width=img_width)
            elif num_images == 1:
                img = data['images'][0]
                img_path = resolve_image_path(img['path'], base_dir, md_dir)
                if img_path:
                    add_image_to_slide(slide, img_path, Inches(3.5), Inches(1.3), width=Inches(6))

            # Extract text from table cells (skip image cells and empty cells)
            text_items = []
            for row in table_data:
                for cell in row:
                    cell_text = strip_markdown_images(str(cell)).strip()
                    if cell_text and cell_text != '|':
                        text_items.append(cell_text)

            # Render text below images
            if text_items:
                text_top = Inches(4.8)
                shape = slide.shapes.add_textbox(
                    Layout.CONTENT_LEFT, text_top,
                    Layout.CONTENT_WIDTH, Inches(2.5)
                )
                tf = shape.text_frame
                tf.word_wrap = True

                first_para = True
                for text in text_items:
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    add_formatted_text(p, text, Fonts.BODY_SIZE, Colors.TEXT_DARK)
                    p.space_after = Pt(6)

            return slide

    # Split content into before-table and after-table based on raw markdown
    raw = data['raw']
    content = data.get('content', [])

    # Find where table starts and ends in raw markdown
    table_start = raw.find('|')
    table_end = raw.rfind('|')

    content_above = []
    content_below = []

    if content and table_start > 0:
        # Get text before and after table
        text_before = raw[:table_start]
        text_after = raw[table_end+1:] if table_end > 0 else ''

        for item_type, text in content:
            # Check if this content appears before or after the table
            if text in text_before:
                content_above.append((item_type, text))
            elif text in text_after:
                content_below.append((item_type, text))

    # Render content above table
    content_top = Inches(1.2)
    if content_above:
        shape = slide.shapes.add_textbox(
            Layout.CONTENT_LEFT, content_top,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        first_para = True
        for item_type, text in content_above:
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            add_formatted_text(p, text, Fonts.BODY_SIZE, Colors.TEXT_DARK)
            p.space_after = Pt(4)
        content_top = Inches(2.0)  # Move table down

    # Table
    if data['table']:
        table_data = data['table']
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 2

        # Use full width and smaller font for wide tables
        table_width = Inches(12.3)
        font_size = Fonts.TABLE_SIZE if cols <= 4 else Fonts.SMALL_SIZE
        row_height = Inches(0.35) if cols <= 4 else Inches(0.3)

        table_shape = slide.shapes.add_table(
            rows, cols,
            Layout.CONTENT_LEFT, content_top,
            table_width, row_height * rows
        )
        table = table_shape.table

        # Calculate column widths based on content
        import math
        col_widths = []
        min_width = Inches(1.0)

        for c_idx in range(cols):
            lengths = []
            for row in table_data:
                if c_idx < len(row):
                    clean = re.sub(r'\*+([^*]+)\*+', r'\1', str(row[c_idx]))
                    clean = re.sub(r'<[^>]+>', '', clean)
                    lengths.append(len(clean))

            if lengths:
                avg_len = sum(lengths) / len(lengths)
                max_len = max(lengths)
                effective_len = math.sqrt(0.7 * avg_len + 0.3 * max_len + 5)
            else:
                effective_len = 3
            col_widths.append(effective_len)

        total_width = sum(col_widths)
        for c_idx in range(cols):
            proportion = col_widths[c_idx] / total_width
            width = int(table_width * proportion)
            table.columns[c_idx].width = max(width, min_width)

        # Fill table
        for r_idx, row in enumerate(table_data):
            is_header = r_idx == 0
            for c_idx, cell_text in enumerate(row):
                if c_idx >= cols:
                    continue
                cell = table.cell(r_idx, c_idx)
                cell.text_frame.word_wrap = True
                para = cell.text_frame.paragraphs[0]

                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = Colors.LIGHT_BLUE
                    clean_text = re.sub(r'\*+([^*]+)\*+', r'\1', cell_text)
                    clean_text = strip_markdown_images(clean_text)
                    para.text = clean_text
                    para.font.size = font_size
                    para.font.name = Fonts.FAMILY
                    para.font.color.rgb = Colors.NAVY
                    para.font.bold = True
                else:
                    # Strip image markdown from table cells (images rendered separately)
                    clean_cell_text = strip_markdown_images(cell_text)
                    if clean_cell_text:
                        add_formatted_text(para, clean_cell_text, font_size, Colors.TEXT_DARK)

        # Calculate where content below should start
        table_bottom = content_top + row_height * rows + Inches(0.2)
    else:
        table_bottom = content_top

    # Render content below table
    if content_below:
        shape = slide.shapes.add_textbox(
            Layout.CONTENT_LEFT, table_bottom,
            Layout.CONTENT_WIDTH, Inches(2.5)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        first_para = True
        for item_type, text in content_below:
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            add_formatted_text(p, text, Fonts.BODY_SIZE, Colors.TEXT_DARK)
            p.space_after = Pt(6)

    return slide


def build_image_slide(prs, data, base_dir, md_dir):
    """Build image-focused slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    if data['headers']:
        title = data['headers'][0][1]
        if data['headers'][0][0] == 1:
            add_h1_with_underline(slide, title)
        else:
            add_h2_with_underline(slide, title, Layout.MARGIN_TOP)

    # Add main image - centered horizontally, below title
    if data['images']:
        img = data['images'][0]
        img_path = resolve_image_path(img['path'], base_dir, md_dir)
        if img_path:
            # Calculate centered position
            img_width = Inches(10)
            img_left = (Layout.WIDTH - img_width) / 2
            add_image_to_slide(
                slide, img_path,
                img_left, Inches(1.5),
                width=img_width
            )
        else:
            print(f"   Warning: Could not resolve image path: {img['path']}")

    # Add any text content (bullets and paragraphs) below the image - centered
    content = data.get('content', [])
    if content:
        shape = slide.shapes.add_textbox(
            Layout.CONTENT_LEFT, Inches(5.8),
            Layout.CONTENT_WIDTH, Inches(1.5)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE  # Vertically center text

        first_para = True
        for item_type, text in content:
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            if item_type == 'header':
                level, header_text = text
                p.text = header_text
                style_paragraph(p, Fonts.H3_SIZE, Colors.ORCHID, bold=True)
                p.space_after = Pt(8)
            elif item_type == 'bullet':
                add_formatted_bullet_text(p, text, Fonts.BODY_SIZE, Colors.TEXT_DARK)
                p.space_after = Pt(8)
            else:
                add_formatted_text(p, text, Fonts.BODY_SIZE, Colors.DARK_GRAY)
                p.space_after = Pt(10)

    return slide


def build_content_slide(prs, data, base_dir, md_dir):
    """Build standard content slide with headers and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title area: top of slide
    title_top = Layout.MARGIN_TOP

    # Add H1/H2 headers at top (H3+ are in content flow)
    for level, text in data['headers']:
        if level == 1:
            add_h1_with_underline(slide, text, title_top)
        elif level == 2:
            add_h2_with_underline(slide, text, title_top)
        break  # Only first header at top

    # Content area: fixed position below title, full width, full remaining height
    content_top = Inches(1.5)  # Below title and underline, with more breathing room
    content_height = Inches(5.5)  # Fill to near bottom

    # Check if slide has images - if so, use narrower text width to avoid overlap
    has_image = bool(data['images'])
    if has_image:
        content_width = Inches(6.0)  # Leave room for image on right
    else:
        content_width = Layout.CONTENT_WIDTH

    # Render ALL content (H3+, paragraphs, bullets) in ONE text box
    content = data.get('content', [])
    if content:
        shape = slide.shapes.add_textbox(
            Layout.CONTENT_LEFT, content_top,
            content_width, content_height
        )
        tf = shape.text_frame
        tf.word_wrap = True
        if has_image:
            # When image present, auto-size text box to avoid overlap
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tf.anchor = MSO_ANCHOR.TOP
        else:
            tf.anchor = MSO_ANCHOR.MIDDLE  # Vertically center text

        first_para = True
        for item_type, text in content:
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            if item_type == 'header':
                level, header_text = text
                p.text = header_text
                # H3+ styled with Orchid to match PDF accent color
                style_paragraph(p, Fonts.H3_SIZE, Colors.ORCHID, bold=True)
                p.space_before = Pt(8)
                p.space_after = Pt(10)

            elif item_type == 'bullet':
                # Use formatted bullet text with lime-colored bullet
                add_formatted_bullet_text(p, text, Fonts.BODY_SIZE, Colors.TEXT_DARK)
                p.space_after = Pt(8)

            elif item_type == 'paragraph':
                # Use formatted text to handle inline bold/italic
                add_formatted_text(p, text, Fonts.BODY_SIZE, Colors.DARK_GRAY)
                p.space_after = Pt(14)  # More space after paragraphs

    # Add image if present (to the right, with gap from text)
    if data['images']:
        img = data['images'][0]
        img_path = resolve_image_path(img['path'], base_dir, md_dir)
        if img_path:
            add_image_to_slide(
                slide, img_path,
                Inches(7.25), Inches(1.5),
                width=Inches(5.5)
            )

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN CONVERTER
# ═══════════════════════════════════════════════════════════════════════════════

def convert_to_pptx(md_file, base_dir, output_path=None):
    """
    Convert Marp markdown to PowerPoint.

    Returns True on success, False on failure.
    """
    # Validate input
    if not os.path.exists(md_file):
        print(f"Error: File not found: {md_file}")
        return False

    print("\n" + "=" * 70)
    print("           CONVERTING TO POWERPOINT")
    print("=" * 70)

    # Read markdown
    print(f"\n   Reading: {os.path.basename(md_file)}")
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Parse markdown
    print("   Parsing slides...")
    slides_data = parse_markdown(content)
    print(f"   Found {len(slides_data)} slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Layout.WIDTH
    prs.slide_height = Layout.HEIGHT

    # Get paths
    md_dir = os.path.dirname(os.path.abspath(md_file))

    # Build each slide
    print("   Building slides...")
    builders = {
        'title': build_title_slide,
        'agenda': build_agenda_slide,
        'break': build_break_slide,
        'section': build_section_slide,
        'two_column': build_two_column_slide,
        'table': build_table_slide,
        'image': build_image_slide,
        'content': build_content_slide,
    }

    type_counts = {}
    for i, data in enumerate(slides_data):
        slide_type = detect_slide_type(data, i)
        type_counts[slide_type] = type_counts.get(slide_type, 0) + 1

        builder = builders.get(slide_type, build_content_slide)
        try:
            builder(prs, data, base_dir, md_dir)
        except Exception as e:
            print(f"   Warning: Error building slide {i+1}: {e}")
            # Build as content slide fallback
            try:
                build_content_slide(prs, data, base_dir, md_dir)
            except:
                pass

    # Show type breakdown
    print("   Slide types:")
    for stype, count in sorted(type_counts.items()):
        print(f"      {stype}: {count}")

    # Save
    if output_path is None:
        output_path = md_file.replace('.md', '.pptx')

    print(f"\n   Saving: {os.path.basename(output_path)}")
    prs.save(output_path)

    file_size = os.path.getsize(output_path) / 1024

    print("\n" + "=" * 70)
    print("                    SUCCESS!")
    print("=" * 70)
    print(f"\n   Output: {output_path}")
    print(f"   Size: {file_size:.1f} KB")
    print(f"   Slides: {len(prs.slides)}")
    print("\n" + "=" * 70 + "\n")

    return True


# ═══════════════════════════════════════════════════════════════════════════════
# CLI INTERFACE
# ═══════════════════════════════════════════════════════════════════════════════

def list_available_decks(base_dir):
    """List markdown decks in outputs/ folder."""
    outputs_dir = os.path.join(base_dir, "outputs")
    if not os.path.exists(outputs_dir):
        return []

    decks = []
    for file in os.listdir(outputs_dir):
        if file.endswith('.md') and not file.startswith('.'):
            decks.append(file)

    return sorted(decks)


def prompt_for_deck(base_dir):
    """Interactive mode: ask user which deck to convert."""
    print("\n" + "=" * 70)
    print("              AVAILABLE DECKS")
    print("=" * 70 + "\n")

    decks = list_available_decks(base_dir)

    if not decks:
        print("No decks found in outputs/ folder!")
        print("\nBuild a deck first:")
        print("   python3 tools/02_build_deck.py --workshop YOUR-WORKSHOP")
        sys.exit(1)

    for i, deck in enumerate(decks, 1):
        deck_path = os.path.join(base_dir, "outputs", deck)
        size = os.path.getsize(deck_path) / 1024
        print(f"  {i}. {deck} ({size:.1f} KB)")

    print("\n" + "-" * 70)

    while True:
        try:
            choice = input("\nWhich deck to convert? (number or name): ").strip()

            if choice.isdigit():
                idx = int(choice) - 1
                if 0 <= idx < len(decks):
                    return decks[idx]

            if choice in decks:
                return choice

            if not choice.endswith('.md'):
                if choice + '.md' in decks:
                    return choice + '.md'

            print(f"Invalid choice. Enter 1-{len(decks)} or a deck name.")

        except KeyboardInterrupt:
            print("\n\nCancelled.")
            sys.exit(0)


def main():
    """Main entry point."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    base_dir = os.path.dirname(script_dir)

    if len(sys.argv) > 1 and not sys.argv[1].startswith('-'):
        # Command line mode
        parser = argparse.ArgumentParser(
            description="Convert Marp markdown to PowerPoint with FASTR styling"
        )
        parser.add_argument('markdown_file', help='Markdown file to convert')
        parser.add_argument('--output', '-o', help='Output PPTX filename')

        args = parser.parse_args()

        md_file = args.markdown_file
        if not md_file.startswith('/'):
            md_file = os.path.join(base_dir, md_file)

        success = convert_to_pptx(md_file, base_dir, args.output)
        sys.exit(0 if success else 1)

    else:
        # Interactive mode
        print("\n" + "=" * 70)
        print("         FASTR POWERPOINT CONVERTER")
        print("=" * 70)

        deck_file = prompt_for_deck(base_dir)
        deck_path = os.path.join(base_dir, "outputs", deck_file)

        success = convert_to_pptx(deck_path, base_dir)
        sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
